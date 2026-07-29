#!/usr/bin/env python3
"""Measure how finely a deck gets split into review fragments.

The review screen lists one item per PPT paragraph (``PPTParser`` →
``ReviewSession.fragments``). Decks that use hard returns for line wrapping,
grouped shapes, or tables therefore scatter a single sentence across several
review items.

This script reports that split empirically, and prototypes the *container path*
that ``ParagraphInfo`` needs in order to group paragraphs back into blocks:
group nesting is currently flattened onto the parent shape index, and table
cells restart ``paragraph_index`` at 0, so ``(slide, shape, paragraph)`` is not
unique today.

Usage:
    python scripts/analyze_fragments.py deck.pptx [--notes] [--json out.json]
    python scripts/analyze_fragments.py deck.pptx --list-blocks
"""

from __future__ import annotations

import argparse
import json
from collections import Counter
from dataclasses import dataclass, field
from pathlib import Path
from typing import Iterator, List, Sequence

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

A_NS = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

# Buckets used for the paragraph-length histogram (upper bound, inclusive).
LENGTH_BUCKETS: Sequence[tuple[str, int]] = (
    ("1-5", 5),
    ("6-10", 10),
    ("11-20", 20),
    ("21-40", 40),
    ("41-80", 80),
    ("81-160", 160),
)


@dataclass
class Fragment:
    """One extracted paragraph, plus the container it belongs to."""

    slide: int
    block_id: str
    block_kind: str  # text_frame | table_cell | notes
    paragraph_index: int
    text: str
    line_breaks: int  # number of <a:br/> elements dropped by runs-join
    runs_text: str  # what PPTParser currently extracts (runs only)


@dataclass
class Block:
    """A container whose paragraphs read as one unit (one text frame/cell)."""

    block_id: str
    slide: int
    kind: str
    fragments: List[Fragment] = field(default_factory=list)

    @property
    def joined(self) -> str:
        return " ".join(f.text for f in self.fragments)


def _count_line_breaks(paragraph) -> int:
    return len(paragraph._p.findall("a:br", A_NS))


def _paragraphs_of(text_frame, slide: int, block_id: str, kind: str) -> Iterator[Fragment]:
    for para_idx, paragraph in enumerate(text_frame.paragraphs):
        runs_text = "".join(run.text for run in paragraph.runs)
        # paragraph.text keeps <a:br/> as "\v" and includes <a:fld> text; the
        # difference against runs_text is exactly what extraction drops today.
        full_text = paragraph.text
        if not runs_text or not runs_text.strip():
            continue
        yield Fragment(
            slide=slide,
            block_id=block_id,
            block_kind=kind,
            paragraph_index=para_idx,
            text=full_text,
            line_breaks=_count_line_breaks(paragraph),
            runs_text=runs_text,
        )


def _walk_shape(shape, slide: int, path: str) -> Iterator[Fragment]:
    """Yield fragments, building a container path that is unique per text frame."""
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child_idx, child in enumerate(shape.shapes):
            yield from _walk_shape(child, slide, f"{path}/g{child_idx}")
        return

    if getattr(shape, "has_table", False):
        for row_idx, row in enumerate(shape.table.rows):
            for col_idx, cell in enumerate(row.cells):
                if not getattr(cell, "text_frame", None):
                    continue
                yield from _paragraphs_of(
                    cell.text_frame, slide, f"{path}/r{row_idx}c{col_idx}", "table_cell"
                )

    if getattr(shape, "has_text_frame", False):
        yield from _paragraphs_of(shape.text_frame, slide, path, "text_frame")


def extract(path: Path, *, notes: bool) -> List[Fragment]:
    presentation = Presentation(str(path))
    fragments: List[Fragment] = []
    for slide_idx, slide in enumerate(presentation.slides, start=1):
        for shape_idx, shape in enumerate(slide.shapes):
            fragments.extend(_walk_shape(shape, slide_idx, f"s{slide_idx}/sh{shape_idx}"))
        if notes and slide.has_notes_slide:
            notes_tf = getattr(slide.notes_slide, "notes_text_frame", None)
            if notes_tf is not None:
                fragments.extend(
                    _paragraphs_of(notes_tf, slide_idx, f"s{slide_idx}/notes", "notes")
                )
    return fragments


def build_blocks(fragments: Sequence[Fragment]) -> List[Block]:
    blocks: dict[str, Block] = {}
    for fragment in fragments:
        block = blocks.get(fragment.block_id)
        if block is None:
            block = Block(fragment.block_id, fragment.slide, fragment.block_kind)
            blocks[fragment.block_id] = block
        block.fragments.append(fragment)
    return list(blocks.values())


def _length_histogram(lengths: Sequence[int]) -> dict[str, int]:
    histogram = {label: 0 for label, _ in LENGTH_BUCKETS}
    histogram["161+"] = 0
    for length in lengths:
        for label, upper in LENGTH_BUCKETS:
            if length <= upper:
                histogram[label] += 1
                break
        else:
            histogram["161+"] += 1
    return histogram


def summarize(fragments: Sequence[Fragment], blocks: Sequence[Block]) -> dict:
    lengths = [len(f.text.strip()) for f in fragments]
    multi = [b for b in blocks if len(b.fragments) > 1]
    kinds = Counter(f.block_kind for f in fragments)
    dropped = [f for f in fragments if f.line_breaks]

    return {
        "fragments": len(fragments),
        "blocks": len(blocks),
        "reduction_pct": round(
            (1 - len(blocks) / len(fragments)) * 100, 1
        ) if fragments else 0.0,
        "fragments_by_kind": dict(kinds),
        "short_fragments": {
            "under_10_chars": sum(1 for n in lengths if n < 10),
            "under_20_chars": sum(1 for n in lengths if n < 20),
            "pct_under_20": round(
                sum(1 for n in lengths if n < 20) / len(lengths) * 100, 1
            ) if lengths else 0.0,
        },
        "length_histogram": _length_histogram(lengths),
        "split_blocks": {
            "count": len(multi),
            "pct_of_blocks": round(len(multi) / len(blocks) * 100, 1) if blocks else 0.0,
            "paragraphs_per_block": dict(
                sorted(Counter(len(b.fragments) for b in blocks).items())
            ),
            "max_paragraphs": max((len(b.fragments) for b in blocks), default=0),
        },
        "line_breaks_dropped": {
            "paragraphs": len(dropped),
            "total_breaks": sum(f.line_breaks for f in dropped),
            "note": "runs-join drops <a:br/>; those lines are glued with no separator",
        },
    }


def print_report(summary: dict, blocks: Sequence[Block], *, list_blocks: bool) -> None:
    frag = summary["fragments"]
    print(f"\n조각(문단) {frag}개 → 블록 {summary['blocks']}개 "
          f"({summary['reduction_pct']}% 감소)")
    print(f"  유형별: {summary['fragments_by_kind']}")

    short = summary["short_fragments"]
    print(f"\n짧은 조각: 20자 미만 {short['under_20_chars']}개 ({short['pct_under_20']}%), "
          f"10자 미만 {short['under_10_chars']}개")
    print(f"  길이 분포: {summary['length_histogram']}")

    split = summary["split_blocks"]
    print(f"\n쪼개진 블록: {split['count']}개 ({split['pct_of_blocks']}% of blocks), "
          f"최대 {split['max_paragraphs']}문단")
    print(f"  블록당 문단 수: {split['paragraphs_per_block']}")

    breaks = summary["line_breaks_dropped"]
    print(f"\n줄바꿈(<a:br/>) 손실: 문단 {breaks['paragraphs']}개, "
          f"총 {breaks['total_breaks']}곳")

    if list_blocks:
        print("\n--- 2문단 이상 블록 ---")
        for block in sorted(
            (b for b in blocks if len(b.fragments) > 1),
            key=lambda b: -len(b.fragments),
        )[:40]:
            print(f"\n[슬라이드 {block.slide}] {block.kind} · {len(block.fragments)}문단")
            for fragment in block.fragments:
                print(f"    {fragment.paragraph_index}: {fragment.text!r}")


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("pptx", type=Path, help="분석할 PPTX 파일")
    parser.add_argument("--notes", action="store_true", help="발표자 노트도 포함")
    parser.add_argument("--list-blocks", action="store_true", help="쪼개진 블록 원문 출력")
    parser.add_argument("--json", type=Path, help="요약을 JSON으로 저장")
    args = parser.parse_args()

    fragments = extract(args.pptx, notes=args.notes)
    if not fragments:
        print("추출된 문단이 없습니다.")
        return

    blocks = build_blocks(fragments)
    summary = summarize(fragments, blocks)
    print_report(summary, blocks, list_blocks=args.list_blocks)

    if args.json:
        args.json.write_text(
            json.dumps(summary, ensure_ascii=False, indent=2), encoding="utf-8"
        )
        print(f"\n저장: {args.json}")


if __name__ == "__main__":
    main()
