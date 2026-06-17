#!/usr/bin/env python3
"""Audit color/style mapping between a source PPTX and translated PPTX.

This script is intentionally heuristic. It extracts paragraph runs and flags
cases that are useful for human review, especially position-like highlights
that can indicate a fallback split rather than semantic color mapping.
"""

from __future__ import annotations

import argparse
import json
from dataclasses import asdict, dataclass
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


A_NS = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
NEUTRAL_SCHEME_COLORS = {
    "tx1",
    "tx2",
    "dk1",
    "dk2",
    "lt1",
    "lt2",
    "bg1",
    "bg2",
}


@dataclass
class RunInfo:
    text: str
    color: str
    neutral: bool
    start: int
    end: int


@dataclass
class ParagraphInfo:
    slide: int
    shape: int
    paragraph: int
    text: str
    runs: list[RunInfo]


@dataclass
class Finding:
    severity: str
    kind: str
    index: int
    location: str
    source_text: str
    target_text: str
    source_runs: list[dict]
    target_runs: list[dict]
    note: str


def iter_shapes(shapes) -> Iterable:
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)
        else:
            yield shape


def is_neutral_rgb(value: str) -> bool:
    if len(value) != 6:
        return False
    try:
        r = int(value[0:2], 16)
        g = int(value[2:4], 16)
        b = int(value[4:6], 16)
    except ValueError:
        return False
    return max(r, g, b) - min(r, g, b) <= 12


def run_color(run) -> tuple[str, bool]:
    rpr = run._r.rPr
    if rpr is None:
        return "default", True

    scheme = rpr.find(".//a:schemeClr", A_NS)
    if scheme is not None:
        value = scheme.get("val") or ""
        return f"scheme:{value}", value in NEUTRAL_SCHEME_COLORS

    srgb = rpr.find(".//a:srgbClr", A_NS)
    if srgb is not None:
        value = (srgb.get("val") or "").upper()
        return f"#{value}", is_neutral_rgb(value)

    return "rPr-no-color", True


def extract_paragraphs(path: Path) -> list[ParagraphInfo]:
    prs = Presentation(str(path))
    paragraphs: list[ParagraphInfo] = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        for shape_index, shape in enumerate(iter_shapes(slide.shapes), start=1):
            if not getattr(shape, "has_text_frame", False):
                continue
            for para_index, paragraph in enumerate(shape.text_frame.paragraphs, start=1):
                cursor = 0
                runs: list[RunInfo] = []
                for run in paragraph.runs:
                    if not run.text:
                        continue
                    color, neutral = run_color(run)
                    start = cursor
                    cursor += len(run.text)
                    runs.append(
                        RunInfo(
                            text=run.text,
                            color=color,
                            neutral=neutral,
                            start=start,
                            end=cursor,
                        )
                    )
                if not runs:
                    continue
                paragraphs.append(
                    ParagraphInfo(
                        slide=slide_index,
                        shape=shape_index,
                        paragraph=para_index,
                        text="".join(run.text for run in runs),
                        runs=runs,
                    )
                )

    return paragraphs


def non_neutral_runs(paragraph: ParagraphInfo) -> list[RunInfo]:
    return [run for run in paragraph.runs if not run.neutral and run.text.strip()]


def relative_position(run: RunInfo, text_len: int) -> str:
    if text_len <= 0:
        return "unknown"
    midpoint = (run.start + run.end) / 2
    ratio = midpoint / text_len
    if ratio < 0.33:
        return "front"
    if ratio > 0.67:
        return "back"
    return "middle"


def location(paragraph: ParagraphInfo) -> str:
    return f"slide {paragraph.slide}, shape {paragraph.shape}, para {paragraph.paragraph}"


def audit(source: list[ParagraphInfo], target: list[ParagraphInfo]) -> list[Finding]:
    findings: list[Finding] = []
    count = min(len(source), len(target))

    for idx in range(count):
        src = source[idx]
        dst = target[idx]
        src_emphasis = non_neutral_runs(src)
        dst_emphasis = non_neutral_runs(dst)

        if not src_emphasis and dst_emphasis:
            findings.append(
                Finding(
                    severity="medium",
                    kind="target_added_highlight",
                    index=idx,
                    location=location(dst),
                    source_text=src.text,
                    target_text=dst.text,
                    source_runs=[asdict(run) for run in src.runs],
                    target_runs=[asdict(run) for run in dst.runs],
                    note="Target has non-neutral highlights where source paragraph had none.",
                )
            )
            continue

        if src_emphasis and not dst_emphasis:
            findings.append(
                Finding(
                    severity="low",
                    kind="dropped_highlight",
                    index=idx,
                    location=location(dst),
                    source_text=src.text,
                    target_text=dst.text,
                    source_runs=[asdict(run) for run in src.runs],
                    target_runs=[asdict(run) for run in dst.runs],
                    note=(
                        "Source highlight was dropped. This is safer than a wrong "
                        "highlight, but may need review."
                    ),
                )
            )
            continue

        if not src_emphasis or not dst_emphasis:
            continue

        src_positions = {
            relative_position(run, len(src.text)) for run in src_emphasis
        }
        dst_positions = {
            relative_position(run, len(dst.text)) for run in dst_emphasis
        }

        if src_positions & dst_positions:
            severity = "high" if len(src.text) != len(dst.text) else "medium"
            findings.append(
                Finding(
                    severity=severity,
                    kind="position_like_highlight",
                    index=idx,
                    location=location(dst),
                    source_text=src.text,
                    target_text=dst.text,
                    source_runs=[asdict(run) for run in src.runs],
                    target_runs=[asdict(run) for run in dst.runs],
                    note=(
                        "Source and target highlights occupy similar relative positions. "
                        "Review for position-based color carryover."
                    ),
                )
            )

    if len(source) != len(target):
        dummy_source = source[-1] if source else ParagraphInfo(0, 0, 0, "", [])
        dummy_target = target[-1] if target else ParagraphInfo(0, 0, 0, "", [])
        findings.append(
            Finding(
                severity="high",
                kind="paragraph_count_mismatch",
                index=count,
                location="document",
                source_text=dummy_source.text,
                target_text=dummy_target.text,
                source_runs=[],
                target_runs=[],
                note=f"Source has {len(source)} paragraphs, target has {len(target)} paragraphs.",
            )
        )

    return findings


def print_text_report(findings: list[Finding], *, limit: int) -> None:
    shown = findings[:limit]
    print(f"Findings: {len(findings)} total, showing {len(shown)}")
    counts: dict[str, int] = {}
    for finding in findings:
        counts[finding.kind] = counts.get(finding.kind, 0) + 1
    for kind, count in sorted(counts.items()):
        print(f"- {kind}: {count}")

    for finding in shown:
        print()
        print(
            f"[{finding.severity}] {finding.kind} @ "
            f"{finding.location} (idx={finding.index})"
        )
        print(f"  source: {finding.source_text}")
        print(f"  target: {finding.target_text}")
        print(f"  note: {finding.note}")
        print("  source runs:")
        for run in finding.source_runs:
            print(f"    {run['color']:14} {run['text']!r}")
        print("  target runs:")
        for run in finding.target_runs:
            print(f"    {run['color']:14} {run['text']!r}")


def main() -> int:
    parser = argparse.ArgumentParser(
        description=(
            "Audit suspicious color/style mapping between source and translated "
            "PPTX files."
        )
    )
    parser.add_argument("source", type=Path, help="Source PPTX")
    parser.add_argument("target", type=Path, help="Translated PPTX")
    parser.add_argument("--json", action="store_true", help="Print JSON findings")
    parser.add_argument(
        "--limit",
        type=int,
        default=30,
        help="Max findings to show in text mode",
    )
    args = parser.parse_args()

    source = extract_paragraphs(args.source)
    target = extract_paragraphs(args.target)
    findings = audit(source, target)

    if args.json:
        print(
            json.dumps(
                [asdict(finding) for finding in findings],
                ensure_ascii=False,
                indent=2,
            )
        )
    else:
        print_text_report(findings, limit=args.limit)

    return 1 if any(f.severity == "high" for f in findings) else 0


if __name__ == "__main__":
    raise SystemExit(main())
