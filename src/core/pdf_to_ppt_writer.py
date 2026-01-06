"""Convert Vision-based semantic OCR results to PowerPoint."""

from __future__ import annotations

import io
import logging
import os
from dataclasses import dataclass
from typing import List, Optional, Tuple

from pptx import Presentation
from pptx.util import Emu, Pt

from src.core.pdf_processor import PageOCRResult, TextBlock
from src.core.pdf_to_ppt_helpers import (
    add_text_box,
    analyze_style_with_gpt,
    blocks_to_layout_rects,
    build_inpainted_background,
    calculate_dynamic_height,
    calculate_layout_metrics,
    crop_image_for_rect,
    encode_image_to_base64,
    find_gaps_from_intervals,
    find_overlaps,
    get_average_color,
    get_luminance,
    get_optimal_text_color,
    get_style_for_type,
    hex_to_rgb,
    is_cjk_char,
    merge_nearby_rects,
    merge_overlaps,
    merge_rect_group,
    parse_alignment,
    xy_cut_merge,
)

LOGGER = logging.getLogger(__name__)


@dataclass
class TextBoxStyle:
    """Style configuration."""
    font_name: str = "맑은 고딕"
    background_color: Optional[Tuple[int, int, int]] = (255, 255, 255) # None means adaptive
    text_color: Optional[Tuple[int, int, int]] = (0, 0, 0)             # None means adaptive
    # Reduced padding for tighter fit to original text
    padding_percent: float = 0.01


@dataclass
class LayoutRect:
    """Rectangle in EMU coordinates for layout calculations."""
    left: int
    top: int
    width: int
    height: int
    block_index: int  # Reference to original TextBlock index
    block_type: str = "body"
    text: str = ""
    
    @property
    def right(self) -> int:
        return self.left + self.width
    
    @property
    def bottom(self) -> int:
        return self.top + self.height
    
    def overlaps_with(self, other: "LayoutRect") -> bool:
        """Check if this rect overlaps with another."""
        # No overlap if one is completely to the left, right, above, or below
        if self.right <= other.left or other.right <= self.left:
            return False
        if self.bottom <= other.top or other.bottom <= self.top:
            return False
        return True
    
    def overlap_area(self, other: "LayoutRect") -> int:
        """Calculate overlap area with another rect."""
        x_overlap = max(0, min(self.right, other.right) - max(self.left, other.left))
        y_overlap = max(0, min(self.bottom, other.bottom) - max(self.top, other.top))
        return x_overlap * y_overlap

    def expand(self, x_ratio: float = 0.05, y_ratio: float = 0.01) -> None:
        """
        [안전장치 1] 비대칭 팽창
        좌우(x)는 넉넉히(5%), 상하(y)는 최소한(1%)으로 늘려서
        윗줄/아랫줄 침범 및 이미지 가림 현상을 방지함
        """
        w_expand = int(self.width * x_ratio)
        h_expand = int(self.height * y_ratio)
        
        self.left -= w_expand
        self.top -= h_expand
        self.width += (w_expand * 2)
        self.height += (h_expand * 2)

    def should_merge_with(self, other: "LayoutRect") -> bool:
        """
        [안전장치 2] 조건부 병합 판단 로직
        """
        # 1. 의미론적 타입이 다르면 병합 금지 (예: 제목과 본문이 합쳐지는 참사 방지)
        if self.block_type != other.block_type:
            return False
            
        # 2. 거리가 너무 멀면 병합 금지 (예: 20pt 이상 떨어져 있으면 별개로 취급)
        # 세로 거리 계산
        y_gap = max(0, other.top - self.bottom) if self.top < other.top else max(0, self.top - other.bottom)
        if y_gap > Pt(20): 
            return False
            
        return True

    def union(self, other: "LayoutRect") -> "LayoutRect":
        """Merge with another rect."""
        new_left = min(self.left, other.left)
        new_top = min(self.top, other.top)
        new_right = max(self.right, other.right)
        new_bottom = max(self.bottom, other.bottom)
        
        # Merge text based on top position
        if self.top <= other.top:
            new_text = f"{self.text}\n{other.text}"
        else:
            new_text = f"{other.text}\n{self.text}"
            
        return LayoutRect(
            left=new_left,
            top=new_top,
            width=new_right - new_left,
            height=new_bottom - new_top,
            block_index=self.block_index, # Keep reference to one
            block_type=self.block_type,
            text=new_text
        )


@dataclass
class DetectedStyle:
    """Style metadata returned from Vision analysis."""
    color_hex: str = "#000000"
    is_bold: bool = False
    font_category: str = "sans-serif"
    alignment: str = "left"


class PDFToPPTWriter:
    """Convert Semantic TextBlocks to PowerPoint with robust layout handling."""

    # Minimum gap between text boxes (in EMU) to prevent visual collision
    MIN_GAP_EMU = int(Pt(3))  # 3pt gap

    def __init__(
        self,
        text_style: Optional[TextBoxStyle] = None,
        openai_api_key: Optional[str] = None,
        vision_model: str = "gpt-5.1",
        enable_inpainting: bool = True,
        enable_style_inference: bool = False,
    ) -> None:
        self._text_style = text_style or TextBoxStyle()
        # Slide dimensions will be set dynamically based on the first page
        self._slide_width = Emu(0)
        self._slide_height = Emu(0)
        self._openai_api_key = openai_api_key or os.environ.get("OPENAI_API_KEY", "")
        self._vision_model = vision_model
        self._enable_inpainting = enable_inpainting
        self._enable_style_inference = enable_style_inference



    def create_presentation(
        self,
        ocr_results: List[PageOCRResult],
        include_background: bool = True,
        resolve_overlaps: bool = True,
    ) -> io.BytesIO:
        """
        Create a PowerPoint presentation from OCR results.
        
        Args:
            ocr_results: List of PageOCRResult objects
            include_background: Whether to include the original PDF page as background
            resolve_overlaps: Whether to resolve overlapping text boxes
        
        Returns:
            BytesIO buffer containing the PPTX file
        """
        prs = Presentation()
        
        if not ocr_results:
            buffer = io.BytesIO()
            prs.save(buffer)
            buffer.seek(0)
            return buffer

        # Determine Master Slide Size from First Page
        first_page = ocr_results[0]
        
        # Convert pixels to EMUs using a standard 96 DPI assumption
        self._slide_width = Emu(first_page.image_width * 9525)
        self._slide_height = Emu(first_page.image_height * 9525)

        prs.slide_width = self._slide_width
        prs.slide_height = self._slide_height
        blank_layout = prs.slide_layouts[6]

        for result in ocr_results:
            slide = prs.slides.add_slide(blank_layout)
            
            # Calculate layout metrics for this specific page
            scale, off_x, off_y = calculate_layout_metrics(
                result.image_width,
                result.image_height,
                self._slide_width,
                self._slide_height,
            )

            cover_original = include_background

            # Convert blocks to layout rects
            layout_rects = blocks_to_layout_rects(
                result.text_blocks, scale, off_x, off_y, self._slide_height, cover_original=cover_original
            )

            # 1) Line-level merge to reduce fragmentation
            layout_rects = merge_nearby_rects(layout_rects)

            # 2) Dynamic grid (XY-cut) merge to respect whitespace and columns
            layout_rects = xy_cut_merge(
                layout_rects,
                page_width=self._slide_width,
                page_height=self._slide_height,
            )

            # Optional background inpainting to remove original text
            background_image = result.image
            if include_background and self._enable_inpainting and layout_rects:
                try:
                    background_image = build_inpainted_background(
                        result.image, layout_rects, scale, off_x, off_y
                    )
                except Exception as exc:
                    LOGGER.warning(
                        "Inpainting failed for page %d, falling back to original background: %s",
                        result.page_number,
                        exc,
                    )
                    background_image = result.image

            # Add background image
            if include_background:
                img_buffer = io.BytesIO()
                background_image.save(img_buffer, format="PNG")
                img_buffer.seek(0)
                
                slide.shapes.add_picture(
                    img_buffer,
                    left=off_x,
                    top=off_y,
                    width=int(result.image_width * scale),
                    height=int(result.image_height * scale),
                )
            
            # Resolve overlaps if enabled
            if resolve_overlaps and cover_original and layout_rects:
                pre_overlaps = find_overlaps(layout_rects)
                if pre_overlaps:
                    LOGGER.info(
                        "Page %d: Found %d overlapping block pairs, resolving...",
                        result.page_number, len(pre_overlaps)
                    )
                
                # Use merge strategy instead of push
                layout_rects = merge_overlaps(layout_rects)
                
                post_overlaps = find_overlaps(layout_rects)
                if post_overlaps:
                    LOGGER.warning(
                        "Page %d: %d overlaps remain after resolution",
                        result.page_number, len(post_overlaps)
                    )
            
            # Add text boxes using resolved layout
            for rect in layout_rects:
                detected_style = None
                if self._enable_style_inference:
                    crop = crop_image_for_rect(result.image, rect, scale, off_x, off_y)
                    if crop:
                        detected_style = analyze_style_with_gpt(
                            crop, rect.text, self._openai_api_key, self._vision_model
                        )

                vision_text_color = None
                if 0 <= rect.block_index < len(result.text_blocks):
                    vision_text_color = result.text_blocks[rect.block_index].text_color

                add_text_box(
                    slide=slide,
                    layout_rect=rect,
                    text_style=self._text_style,
                    slide_width=self._slide_width,
                    slide_height=self._slide_height,
                    image=background_image,  # Use cleaned background for adaptive color
                    scale=scale,
                    offset_x=off_x,
                    offset_y=off_y,
                    detected_style=detected_style,
                    vision_text_color=vision_text_color,
                    cover_original=cover_original,
                    enable_inpainting=self._enable_inpainting,
                )

        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer
