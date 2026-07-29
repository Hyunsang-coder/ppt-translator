"""Utilities for parsing PPT/PPTX files into structured paragraph data."""

from __future__ import annotations

import io
import logging
from dataclasses import dataclass
from typing import TYPE_CHECKING, Any, List, Sequence, Tuple

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

if TYPE_CHECKING:  # pragma: no cover - import for type checking only
    from pptx.text.text import Paragraph
else:
    Paragraph = Any

LOGGER = logging.getLogger(__name__)


@dataclass
class ParagraphInfo:
    """Metadata for a single paragraph in the presentation."""

    slide_index: int
    shape_index: int
    paragraph_index: int
    original_text: str
    paragraph: Paragraph
    slide_title: str | None
    is_note: bool = False
    # Identifies the text frame this paragraph lives in. ``shape_index`` alone
    # cannot: grouped shapes flatten onto the parent's index and every cell of a
    # table shares it while restarting ``paragraph_index`` at 0, so
    # (slide, shape, paragraph) is not unique. Consumers that need to know which
    # paragraphs read together — the review screen groups a hard-return-wrapped
    # sentence back into one item — key on this instead.
    container_id: str = ""
    # What kind of text frame it is. A body placeholder with several paragraphs
    # is a bullet list; a plain text box with several is usually one sentence
    # the author wrapped by hand. Bullet markers themselves are inherited from
    # the layout's list style and are absent from the paragraph XML, so the
    # placeholder type is the reliable signal.
    container_kind: str = "textbox"


_TITLE_PLACEHOLDERS = {"TITLE", "CENTER_TITLE", "VERTICAL_TITLE"}
_BODY_PLACEHOLDERS = {"BODY", "SUBTITLE", "OBJECT", "VERTICAL_BODY", "VERTICAL_OBJECT"}


def _container_kind(shape) -> str:
    """Classify a text frame by what it is on the slide.

    Used to tell a bullet list (body placeholder) from a sentence the author
    wrapped with hard returns (plain text box), which the paragraph XML cannot
    distinguish on its own.
    """
    try:
        if not shape.is_placeholder:
            return "textbox"
        name = str(shape.placeholder_format.type).split()[0].upper()
    except (AttributeError, ValueError):  # pragma: no cover - odd shapes
        return "textbox"
    if name in _TITLE_PLACEHOLDERS:
        return "title"
    if name in _BODY_PLACEHOLDERS:
        return "body"
    return "placeholder"


class PPTParser:
    """Extract text content and structure from PPT presentations."""

    def extract_paragraphs(
        self, ppt_file: io.BytesIO, *, translate_notes: bool = False
    ) -> Tuple[List[ParagraphInfo], Presentation]:
        """Parse the uploaded PPT file and collect paragraphs.

        Args:
            ppt_file: In-memory PPT/PPTX file buffer.
            translate_notes: If True, also extract speaker notes paragraphs.

        Returns:
            A tuple containing the list of paragraphs and the loaded presentation object.
        """

        ppt_file.seek(0)
        presentation = Presentation(ppt_file)
        paragraphs: List[ParagraphInfo] = []

        for slide_index, slide in enumerate(presentation.slides):
            slide_title = None
            if slide.shapes.title and slide.shapes.title.text:
                slide_title = slide.shapes.title.text

            for shape_index, shape in enumerate(slide.shapes):
                paragraphs.extend(
                    self._extract_from_shape(
                        shape=shape,
                        slide_index=slide_index,
                        shape_index=shape_index,
                        slide_title=slide_title,
                        container_path=f"s{slide_index}/sh{shape_index}",
                    )
                )

            if translate_notes and slide.has_notes_slide:
                notes_slide = slide.notes_slide
                notes_tf = getattr(notes_slide, "notes_text_frame", None)
                if notes_tf is not None:
                    for para_idx, paragraph in enumerate(notes_tf.paragraphs):
                        text = "".join(run.text for run in paragraph.runs)
                        if not text or not text.strip():
                            continue
                        paragraphs.append(
                            ParagraphInfo(
                                slide_index=slide_index,
                                shape_index=-1,
                                paragraph_index=para_idx,
                                original_text=text,
                                paragraph=paragraph,
                                slide_title=slide_title,
                                is_note=True,
                                container_id=f"s{slide_index}/notes",
                                container_kind="notes",
                            )
                        )

        LOGGER.info("Extracted %d paragraphs from %d slides.", len(paragraphs), len(presentation.slides))
        return paragraphs, presentation

    def _extract_from_shape(
        self,
        shape,
        slide_index: int,
        shape_index: int,
        slide_title: str | None,
        container_path: str = "",
    ) -> List[ParagraphInfo]:
        """Recursively extract paragraphs from a shape and its children."""

        collected: List[ParagraphInfo] = []

        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            # Children keep the parent's shape_index (translation and writing key
            # on the paragraph objects, not the index), so the path is what keeps
            # each child's text frame distinguishable.
            for child_index, child in enumerate(shape.shapes):  # type: ignore[attr-defined]
                collected.extend(
                    self._extract_from_shape(
                        shape=child,
                        slide_index=slide_index,
                        shape_index=shape_index,
                        slide_title=slide_title,
                        container_path=f"{container_path}/g{child_index}",
                    )
                )
            return collected

        if getattr(shape, "has_table", False):
            table = shape.table
            collected.extend(
                self._extract_from_table(
                    table=table,
                    slide_index=slide_index,
                    shape_index=shape_index,
                    slide_title=slide_title,
                    container_path=container_path,
                )
            )

        if getattr(shape, "has_text_frame", False):
            text_frame = shape.text_frame
            collected.extend(
                self._collect_paragraphs_from_text_frame(
                    text_frame.paragraphs,
                    slide_index,
                    shape_index,
                    slide_title,
                    container_id=container_path,
                    container_kind=_container_kind(shape),
                )
            )

        return collected

    def _extract_from_table(
        self,
        table,
        slide_index: int,
        shape_index: int,
        slide_title: str | None,
        container_path: str = "",
    ) -> List[ParagraphInfo]:
        """Extract paragraphs from a table shape."""

        collected: List[ParagraphInfo] = []
        for row_index, row in enumerate(table.rows):
            for column_index, cell in enumerate(row.cells):
                if not getattr(cell, "text_frame", None):
                    continue
                collected.extend(
                    self._collect_paragraphs_from_text_frame(
                        cell.text_frame.paragraphs,
                        slide_index,
                        shape_index,
                        slide_title,
                        # Every cell shares the table's shape_index and restarts
                        # paragraph_index at 0; the cell coordinates are what
                        # make each cell's paragraphs addressable.
                        container_id=f"{container_path}/r{row_index}c{column_index}",
                        container_kind="table_cell",
                    )
                )
        return collected

    @staticmethod
    def _collect_paragraphs_from_text_frame(
        paragraphs: Sequence[Paragraph],
        slide_index: int,
        shape_index: int,
        slide_title: str | None,
        container_id: str = "",
        container_kind: str = "textbox",
    ) -> List[ParagraphInfo]:
        """Convert pptx paragraphs into ParagraphInfo instances."""

        collected: List[ParagraphInfo] = []
        for paragraph_index, paragraph in enumerate(paragraphs):
            text = "".join(run.text for run in paragraph.runs)
            if not text or not text.strip():
                continue

            collected.append(
                ParagraphInfo(
                    slide_index=slide_index,
                    shape_index=shape_index,
                    paragraph_index=paragraph_index,
                    original_text=text,
                    paragraph=paragraph,
                    slide_title=slide_title,
                    container_id=container_id,
                    container_kind=container_kind,
                )
            )
        return collected
