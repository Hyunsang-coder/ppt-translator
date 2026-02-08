"""Write translated text back into the original PPT presentation."""

from __future__ import annotations

import io
import logging
from typing import Iterable, List

from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt

from src.utils.helpers import split_text_into_segments

LOGGER = logging.getLogger(__name__)

# Expansion threshold: text must grow more than this % to trigger fitting
_EXPANSION_THRESHOLD = 1.10  # 10%

# Valid mode strings (matches TextFitMode enum values)
_MODE_NONE = "none"
_MODE_AUTO_SHRINK = "auto_shrink"
_MODE_EXPAND_BOX = "expand_box"
_MODE_SHRINK_THEN_EXPAND = "shrink_then_expand"

# Safety gap for width expansion (EMU) — approximately 2mm
_EXPANSION_GAP_EMU = 72000


def _build_shape_context(presentation):
    """Build mappings from txBody to shape and per-slide bounding boxes.

    Returns:
        Tuple of (txbody_to_shape, slide_bounds, slide_width) where:
        - txbody_to_shape: dict mapping id(txBody) -> (shape, slide_idx)
        - slide_bounds: dict mapping slide_idx -> list of (left, top, right, bottom, txbody_id)
        - slide_width: presentation slide width in EMU
    """
    txbody_to_shape: dict[int, tuple] = {}
    slide_bounds: dict[int, list] = {}
    slide_width = presentation.slide_width

    for slide_idx, slide in enumerate(presentation.slides):
        bounds_list: list[tuple] = []
        for shape in slide.shapes:
            left = getattr(shape, "left", None)
            top = getattr(shape, "top", None)
            width = getattr(shape, "width", None)
            height = getattr(shape, "height", None)
            if left is None or top is None or width is None or height is None:
                continue

            txbody_id = None
            if getattr(shape, "has_text_frame", False):
                try:
                    txbody_id = id(shape.text_frame._txBody)
                    txbody_to_shape[txbody_id] = (shape, slide_idx)
                except Exception:
                    pass

            bounds_list.append((left, top, left + width, top + height, txbody_id))
        slide_bounds[slide_idx] = bounds_list

    return txbody_to_shape, slide_bounds, slide_width


def _calculate_available_expansion(shape, slide_bounds_list, slide_width, own_txbody_id):
    """Calculate how far a shape can expand left and right (EMU).

    Only shapes with vertical overlap are treated as obstacles.
    """
    s_left = shape.left
    s_top = shape.top
    s_right = s_left + shape.width
    s_bottom = s_top + shape.height

    avail_left = s_left
    avail_right = slide_width - s_right

    for ob_left, ob_top, ob_right, ob_bottom, ob_id in slide_bounds_list:
        if ob_id is not None and ob_id == own_txbody_id:
            continue
        # Skip if no vertical overlap
        if ob_bottom <= s_top or ob_top >= s_bottom:
            continue
        # Obstacle to the left
        if ob_right <= s_left:
            avail_left = min(avail_left, s_left - ob_right)
        # Obstacle to the right
        if ob_left >= s_right:
            avail_right = min(avail_right, ob_left - s_right)

    avail_left = max(0, avail_left - _EXPANSION_GAP_EMU)
    avail_right = max(0, avail_right - _EXPANSION_GAP_EMU)
    return avail_left, avail_right


def _safe_expand_width(shape, available_right, needed_expansion_emu):
    """Expand shape width to the right within available space.

    The shape's left position is preserved so the textbox doesn't shift.
    Returns the width expansion ratio (new_width / old_width, >= 1.0).
    """
    if available_right <= 0 or needed_expansion_emu <= 0:
        return 1.0

    actual_expansion = min(needed_expansion_emu, available_right)

    old_width = shape.width
    shape.width = old_width + actual_expansion

    return shape.width / old_width if old_width > 0 else 1.0


def _update_slide_bounds(slide_bounds, slide_idx, txbody_id, shape):
    """Update the cached bounds for a shape after width expansion."""
    bounds = slide_bounds.get(slide_idx, [])
    for i, entry in enumerate(bounds):
        if entry[4] is not None and entry[4] == txbody_id:
            bounds[i] = (
                shape.left,
                shape.top,
                shape.left + shape.width,
                shape.top + shape.height,
                txbody_id,
            )
            break


def apply_text_fit(
    text_frame,
    original_len: int,
    translated_len: int,
    mode="none",
    min_font_ratio: int = 80,
) -> None:
    """Apply text fitting adjustments to a text frame after translation.

    Args:
        text_frame: python-pptx TextFrame object.
        original_len: Character length of the original text.
        translated_len: Character length of the translated text.
        mode: Text fitting strategy. Accepts TextFitMode enum or string
              ("none", "auto_shrink", "expand_box", "shrink_then_expand").
        min_font_ratio: Minimum font size as percentage of original (50-100).
    """
    # Normalize mode to string value
    mode_str = mode.value if hasattr(mode, "value") else str(mode)

    if mode_str == _MODE_NONE:
        return

    if original_len <= 0 or translated_len <= original_len:
        return

    expansion_ratio = translated_len / original_len
    if expansion_ratio <= _EXPANSION_THRESHOLD:
        return

    if mode_str == _MODE_AUTO_SHRINK:
        _apply_auto_shrink(text_frame, expansion_ratio, min_font_ratio)
    elif mode_str == _MODE_EXPAND_BOX:
        _apply_expand_box(text_frame)
    elif mode_str == _MODE_SHRINK_THEN_EXPAND:
        _apply_shrink_then_expand(text_frame, expansion_ratio, min_font_ratio)


def _apply_auto_shrink(
    text_frame,
    expansion_ratio: float,
    min_font_ratio: int,
) -> None:
    """Shrink fonts proportionally, with a floor ratio and TEXT_TO_FIT_SHAPE fallback."""
    shrink_factor = 1.0 / expansion_ratio
    floor_factor = min_font_ratio / 100.0
    hit_floor = shrink_factor < floor_factor
    effective_factor = max(shrink_factor, floor_factor)

    # Enable word wrap
    text_frame.word_wrap = True

    # Shrink all runs with explicit font sizes
    _shrink_runs(text_frame, effective_factor)

    # If we hit the floor, set TEXT_TO_FIT_SHAPE as a safety net
    if hit_floor:
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


def _apply_expand_box(text_frame) -> None:
    """Expand the text box to fit translated text."""
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT


def _apply_shrink_then_expand(
    text_frame,
    expansion_ratio: float,
    min_font_ratio: int,
) -> None:
    """Shrink fonts first, then expand box if shrink alone isn't enough."""
    shrink_factor = 1.0 / expansion_ratio
    floor_factor = min_font_ratio / 100.0
    hit_floor = shrink_factor < floor_factor
    effective_factor = max(shrink_factor, floor_factor)

    text_frame.word_wrap = True
    _shrink_runs(text_frame, effective_factor)

    # If shrinking wasn't enough, expand the box instead of TEXT_TO_FIT_SHAPE
    if hit_floor:
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT


def _shrink_runs(text_frame, factor: float) -> None:
    """Shrink all runs with explicit font sizes by the given factor."""
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size is not None:
                run.font.size = int(run.font.size * factor)


class PPTWriter:
    """Apply translated strings back into PPT paragraphs."""

    def apply_translations(
        self,
        paragraphs,
        translations: Iterable[str],
        presentation: Presentation,
        text_fit_mode: str = "none",
        min_font_ratio: int = 80,
    ) -> io.BytesIO:
        """Apply translated paragraphs and return the updated PPT as bytes.

        Args:
            paragraphs: Iterable of :class:`ParagraphInfo` objects sharing paragraph references.
            translations: Translated strings aligned to the provided paragraphs.
            presentation: Loaded presentation object to mutate in place.
            text_fit_mode: Text fitting mode ("none", "auto_shrink", "expand_box").
            min_font_ratio: Minimum font size as percentage of original (50-100).

        Returns:
            BytesIO buffer containing the updated PPTX file.
        """

        translation_list: List[str] = list(translations)
        total = len(translation_list)
        fit_adjusted_count = 0
        LOGGER.info("Starting to apply %d translated paragraphs to presentation.", total)

        # Collect text frames that need fitting (deduplicate by text frame id)
        text_frames_to_fit: dict[int, tuple] = {}

        for idx, (paragraph_info, translation) in enumerate(zip(paragraphs, translation_list), start=1):
            if idx % 100 == 0 or idx == total:
                LOGGER.info("Applying translation %d/%d...", idx, total)

            paragraph = paragraph_info.paragraph
            original_text = paragraph_info.original_text
            runs = list(paragraph.runs)

            if not runs:
                run = paragraph.add_run()
                runs = [run]

            weights = [max(len(run.text), 1) for run in runs]
            segments = split_text_into_segments(translation, len(runs), weights=weights)

            for run, segment in zip(runs, segments):
                run.text = segment

            # Track text frames for fitting
            if text_fit_mode != _MODE_NONE:
                text_frame = paragraph._parent
                tf_id = id(text_frame)
                if tf_id not in text_frames_to_fit:
                    text_frames_to_fit[tf_id] = (text_frame, len(original_text), len(translation))
                else:
                    # Accumulate lengths for multi-paragraph text frames
                    _, prev_orig, prev_trans = text_frames_to_fit[tf_id]
                    text_frames_to_fit[tf_id] = (
                        text_frame,
                        prev_orig + len(original_text),
                        prev_trans + len(translation),
                    )

        # Build shape context for width expansion (all non-NONE modes)
        txbody_to_shape: dict = {}
        slide_bounds_map: dict = {}
        slide_w = 0
        if text_fit_mode != _MODE_NONE:
            txbody_to_shape, slide_bounds_map, slide_w = _build_shape_context(presentation)

        # Apply width expansion + text fitting per text frame
        for tf, orig_len, trans_len in text_frames_to_fit.values():
            if trans_len > orig_len:
                effective_orig = orig_len

                # Width expansion before text fit
                if txbody_to_shape and orig_len > 0:
                    exp_ratio = trans_len / orig_len
                    if exp_ratio > _EXPANSION_THRESHOLD:
                        bridge_key = id(tf._txBody)
                        shape_info = txbody_to_shape.get(bridge_key)
                        if shape_info is not None:
                            shape, s_idx = shape_info
                            rotation = getattr(shape, "rotation", 0.0) or 0.0
                            if abs(rotation) < 0.01:
                                needed = int(shape.width * (exp_ratio - 1))
                                bounds = slide_bounds_map.get(s_idx, [])
                                _, ar = _calculate_available_expansion(
                                    shape, bounds, slide_w, bridge_key
                                )
                                w_ratio = _safe_expand_width(shape, ar, needed)
                                if w_ratio > 1.0:
                                    effective_orig = int(orig_len * w_ratio)
                                    _update_slide_bounds(
                                        slide_bounds_map, s_idx, bridge_key, shape
                                    )

                apply_text_fit(
                    tf, effective_orig, trans_len,
                    mode=text_fit_mode, min_font_ratio=min_font_ratio,
                )
                fit_adjusted_count += 1

        if fit_adjusted_count > 0:
            LOGGER.info(
                "Text fitting applied to %d text frames (mode=%s, min_ratio=%d%%).",
                fit_adjusted_count, text_fit_mode, min_font_ratio,
            )

        LOGGER.info("Saving translated presentation to buffer...")
        buffer = io.BytesIO()
        presentation.save(buffer)
        buffer.seek(0)
        LOGGER.info("Applied %d translated paragraphs to presentation.", len(translation_list))
        return buffer
