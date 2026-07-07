"""Utilities for extracting structured text content from PPT presentations."""

from __future__ import annotations

from dataclasses import dataclass, field
from typing import Iterable, List, Literal, Optional, Sequence, Tuple, Union

from pptx import Presentation
from pptx.chart.chart import Chart
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.autoshape import Shape
from pptx.shapes.group import GroupShape
from pptx.shapes.picture import Picture
from pptx.shapes.shapetree import SlideShapes


@dataclass(slots=True)
class TextBlock:
    """Represents a chunk of text extracted from a shape."""

    shape_id: str
    lines: List[str]
    indent_levels: List[int]


@dataclass(slots=True)
class TableBlock:
    """Represents a table with rows of cell text."""

    shape_id: str
    rows: List[List[str]]
    has_header: bool


@dataclass(slots=True)
class FigureBlock:
    """Placeholder metadata for figures and charts."""

    shape_id: str
    figure_type: Literal["image", "chart"]
    title: Optional[str]


@dataclass(slots=True)
class ChartBlock:
    """Chart with extracted category/series data."""

    shape_id: str
    title: Optional[str]
    categories: List[str]
    series: List[Tuple[str, List[Optional[float]]]]


@dataclass(slots=True)
class NoteBlock:
    """Speaker note captured from the slide."""

    text: str


SlideBlock = Union[TextBlock, TableBlock, FigureBlock, ChartBlock, NoteBlock]


@dataclass(slots=True)
class SlideDoc:
    """Structured representation of a single slide."""

    slide_index: int
    title: str
    blocks: List[SlideBlock] = field(default_factory=list)


@dataclass(slots=True)
class ExtractionOptions:
    """Configuration flags controlling extraction behaviour."""

    figures: Literal["omit", "placeholder"] = "omit"
    charts: Literal["labels", "placeholder", "omit"] = "labels"
    table_header: bool = True
    with_notes: bool = False
    slide_range: Optional[Iterable[int]] = None


def _shape_sort_key(shape):
    """Sort key for visual reading order: top-to-bottom, then left-to-right.

    Shapes without position information sort last, preserving relative order.
    """
    try:
        top = shape.top
    except Exception:
        top = None
    try:
        left = shape.left
    except Exception:
        left = None
    return (top is None, top if top is not None else 0, left is None, left if left is not None else 0)


def _iter_shapes(shapes: SlideShapes):
    """Iterate shapes in visual reading order rather than z-order.

    Group shapes are ordered by the group's own position; their children are
    recursively sorted within the group.
    """
    for shape in sorted(shapes, key=_shape_sort_key):
        if isinstance(shape, GroupShape):
            for sub_shape in _iter_shapes(shape.shapes):
                yield sub_shape
        else:
            yield shape


def _shape_text_lines(shape: Shape) -> List[str]:
    lines: List[str] = []
    if not hasattr(shape, "text_frame") or shape.text_frame is None:
        return lines
    for paragraph in shape.text_frame.paragraphs:
        text = "".join(run.text for run in paragraph.runs)
        lines.append(text)
    return lines


def _shape_indent_levels(shape: Shape) -> List[int]:
    levels: List[int] = []
    if not hasattr(shape, "text_frame") or shape.text_frame is None:
        return levels
    for paragraph in shape.text_frame.paragraphs:
        levels.append(paragraph.level or 0)
    return levels


def _sanitize_line(line: str) -> str:
    if line is None:
        return ""
    return line.strip()


def extract_slide(prs_slide, slide_index: int, options: ExtractionOptions) -> SlideDoc:
    """Extract structured content from a single slide."""

    title = None
    title_shape_id: Optional[str] = None
    title_uses_full_shape = False
    title_shape = getattr(prs_slide.shapes, "title", None)
    if title_shape is not None and getattr(title_shape, "has_text_frame", False):
        title = (title_shape.text or "").strip() or None
        if title is not None:
            title_shape_id = str(title_shape.shape_id)
            title_uses_full_shape = True

    if title is None:
        # Fallback: topmost text box in visual reading order (only its first
        # non-empty line is consumed as the title).
        for shape in _iter_shapes(prs_slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX and getattr(shape, "text_frame", None) is not None:
                lines = _shape_text_lines(shape)
                for line in lines:
                    candidate = _sanitize_line(line)
                    if candidate:
                        title = candidate
                        title_shape_id = str(shape.shape_id)
                        break
                if title:
                    break

    if title is None:
        title = f"Slide {slide_index + 1}"

    slide_doc = SlideDoc(slide_index=slide_index, title=title, blocks=[])

    for shape in _iter_shapes(prs_slide.shapes):
        try:
            if title_shape_id is not None and str(shape.shape_id) == title_shape_id:
                if title_uses_full_shape:
                    continue
                # Fallback title: drop the consumed first line, keep the rest.
                lines = _shape_text_lines(shape)
                levels = _shape_indent_levels(shape)
                first_idx = next(
                    (i for i, line in enumerate(lines) if _sanitize_line(line)), None
                )
                if first_idx is not None:
                    lines = lines[first_idx + 1 :]
                    levels = levels[first_idx + 1 :]
                if lines and any(_sanitize_line(line) for line in lines):
                    slide_doc.blocks.append(
                        TextBlock(
                            shape_id=str(shape.shape_id),
                            lines=[_sanitize_line(line) for line in lines if line is not None],
                            indent_levels=levels,
                        )
                    )
                continue
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE and getattr(shape, "table", None):
                rows = []
                for row in shape.table.rows:
                    row_cells = []
                    for cell in row.cells:
                        row_cells.append((cell.text or "").strip())
                    rows.append(row_cells)
                slide_doc.blocks.append(
                    TableBlock(shape_id=str(shape.shape_id), rows=rows, has_header=options.table_header)
                )
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE and isinstance(shape, Picture):
                if options.figures != "omit":
                    slide_doc.blocks.append(
                        FigureBlock(
                            shape_id=str(shape.shape_id),
                            figure_type="image",
                            title=getattr(shape, "name", None),
                        )
                    )
            elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                if options.charts == "labels":
                    chart: Chart = shape.chart  # type: ignore[attr-defined]
                    chart_title = None
                    if chart.has_title:
                        chart_title = chart.chart_title.text_frame.text
                    categories: List[str] = []
                    series: List[Tuple[str, List[Optional[float]]]] = []
                    try:
                        plot = chart.plots[0]
                        categories = [str(c) for c in plot.categories]
                        for s in plot.series:
                            series.append((str(s.name or ""), list(s.values)))
                    except Exception:
                        categories, series = [], []
                    if categories and series:
                        slide_doc.blocks.append(
                            ChartBlock(
                                shape_id=str(shape.shape_id),
                                title=chart_title,
                                categories=categories,
                                series=series,
                            )
                        )
                    else:
                        # Fall back to the previous placeholder behaviour.
                        slide_doc.blocks.append(
                            FigureBlock(
                                shape_id=str(shape.shape_id),
                                figure_type="chart",
                                title=chart_title,
                            )
                        )
                elif options.charts == "placeholder":
                    slide_doc.blocks.append(
                        FigureBlock(
                            shape_id=str(shape.shape_id),
                            figure_type="chart",
                            title=None,
                        )
                    )
            else:
                if getattr(shape, "has_text_frame", False) and shape.text_frame is not None:
                    lines = _shape_text_lines(shape)
                    if lines and any(line.strip() for line in lines):
                        slide_doc.blocks.append(
                            TextBlock(
                                shape_id=str(shape.shape_id),
                                lines=[_sanitize_line(line) for line in lines if line is not None],
                                indent_levels=_shape_indent_levels(shape),
                            )
                        )
        except Exception:  # pragma: no cover - skip problematic shapes without failing entirely
            continue

    if (
        options.with_notes
        and hasattr(prs_slide, "notes_slide")
        and prs_slide.notes_slide
        and getattr(prs_slide.notes_slide, "notes_text_frame", None)
    ):
        note_text = prs_slide.notes_slide.notes_text_frame.text
        if note_text and note_text.strip():
            slide_doc.blocks.append(NoteBlock(text=note_text.strip()))

    return slide_doc


def extract_pptx_to_docs(path_or_buffer, options: ExtractionOptions) -> List[SlideDoc]:
    """Extract structured slide documents from a PPTX path or file-like object."""

    presentation = Presentation(path_or_buffer)
    docs: List[SlideDoc] = []
    slide_filter = set(options.slide_range) if options.slide_range else None

    for idx, slide in enumerate(presentation.slides):
        if slide_filter is not None and (idx + 1) not in slide_filter:
            continue
        docs.append(extract_slide(slide, idx, options))

    return docs


def _md_escape(text: str) -> str:
    normalized = text.replace("\r\n", "\n").replace("\r", "\n")
    return normalized.replace("|", "\\|").replace("\n", "<br>")


def _format_chart_value(value) -> str:
    if value is None:
        return ""
    if isinstance(value, float) and value.is_integer():
        return str(int(value))
    return str(value)


def blocks_to_markdown(blocks: Sequence[SlideBlock], options: ExtractionOptions) -> str:
    lines: List[str] = []
    for block in blocks:
        if isinstance(block, TextBlock):
            for idx, line in enumerate(block.lines):
                if not line:
                    continue
                level = 0
                if block.indent_levels and idx < len(block.indent_levels):
                    level = block.indent_levels[idx]
                lines.append(f"{'  ' * level}- {line}")
            lines.append("")
        elif isinstance(block, TableBlock):
            if not block.rows:
                continue
            header = block.rows[0] if block.has_header and block.rows else None
            body = block.rows[1:] if header else block.rows
            if header:
                lines.append("| " + " | ".join(_md_escape(cell) for cell in header) + " |")
                lines.append("| " + " | ".join(["---"] * len(header)) + " |")
            for row in body:
                lines.append("| " + " | ".join(_md_escape(cell) for cell in row) + " |")
            lines.append("")
        elif isinstance(block, ChartBlock):
            lines.append(f"[Chart: {block.title}]" if block.title else "[Chart]")
            header = ["구분"] + [str(c) for c in block.categories]
            lines.append("| " + " | ".join(_md_escape(cell) for cell in header) + " |")
            lines.append("| " + " | ".join(["---"] * len(header)) + " |")
            width = len(block.categories)
            for name, values in block.series:
                cells = [_format_chart_value(v) for v in values[:width]]
                cells += [""] * (width - len(cells))
                lines.append("| " + " | ".join(_md_escape(cell) for cell in [name, *cells]) + " |")
            lines.append("")
        elif isinstance(block, FigureBlock):
            if block.figure_type == "image" and options.figures == "placeholder":
                title = block.title or "Image"
                lines.append(f"[Figure: {title}]")
            elif block.figure_type == "chart":
                if options.charts == "labels":
                    title = block.title or "Chart"
                    lines.append(f"[Figure: Chart, title=\"{title}\"]")
                elif options.charts == "placeholder":
                    lines.append("[Figure: Chart]")
            lines.append("")
        elif isinstance(block, NoteBlock):
            lines.append("> NOTE: " + block.text.replace("\n", " ").strip())
            lines.append("")
    return "\n".join(lines).strip() + "\n"


def docs_to_markdown(docs: Sequence[SlideDoc], options: ExtractionOptions) -> str:
    chunks: List[str] = []
    for doc in docs:
        chunks.append(f"## Slide {doc.slide_index + 1} - {doc.title}")
        chunks.append("")
        chunks.append(blocks_to_markdown(doc.blocks, options).rstrip())
        chunks.append("")
    return "\n".join(chunks).rstrip() + "\n"


__all__ = [
    "TextBlock",
    "TableBlock",
    "FigureBlock",
    "ChartBlock",
    "NoteBlock",
    "SlideBlock",
    "SlideDoc",
    "ExtractionOptions",
    "extract_slide",
    "extract_pptx_to_docs",
    "blocks_to_markdown",
    "docs_to_markdown",
]
