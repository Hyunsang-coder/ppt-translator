"""Verification loop for multi-color PPT text matching.

This module builds a small multi-color PPTX fixture, runs the same color
distribution path used by translation jobs, writes the translated deck, then
re-opens the output and verifies that important substrings kept their intended
source formatting colors.
"""

from __future__ import annotations

import asyncio
import tempfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Iterable, Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from src.core.ppt_parser import PPTParser
from src.core.ppt_writer import PPTWriter
from src.services.models import DEFAULT_LIGHT_MODEL
from src.services.translation_service import TranslationService


@dataclass(frozen=True, slots=True)
class SourceSegment:
    """One source run with an expected color."""

    text: str
    color: str


@dataclass(frozen=True, slots=True)
class ColorAnchor:
    """A translated substring that should carry a specific source color."""

    text: str
    color: str


@dataclass(frozen=True, slots=True)
class ColorMatchCase:
    """One color matching fixture case."""

    name: str
    source_segments: tuple[SourceSegment, ...]
    translation: str
    anchors: tuple[ColorAnchor, ...]


@dataclass(slots=True)
class ColorMatchCaseResult:
    """Verification result for one case."""

    name: str
    passed: bool
    output_text: str
    errors: list[str] = field(default_factory=list)


@dataclass(slots=True)
class ColorMatchIterationResult:
    """Verification result for one iteration."""

    provider: str
    model: str
    iteration: int
    passed: bool
    cases: list[ColorMatchCaseResult]
    input_path: Path | None = None
    output_path: Path | None = None


DEFAULT_COLOR_MATCH_CASES: tuple[ColorMatchCase, ...] = (
    ColorMatchCase(
        name="numeric-red-anchor",
        source_segments=(
            SourceSegment("Revenue increased by ", "111111"),
            SourceSegment("20%", "D72638"),
            SourceSegment(" in North America", "111111"),
        ),
        translation="북미 지역 매출이 20% 증가했습니다",
        anchors=(ColorAnchor("20%", "D72638"),),
    ),
    ColorMatchCase(
        name="link-blue-reordered",
        source_segments=(
            SourceSegment("Click ", "111111"),
            SourceSegment("here", "1F6FEB"),
            SourceSegment(" to continue", "111111"),
        ),
        translation="계속하려면 여기를 클릭하세요",
        anchors=(ColorAnchor("여기", "1F6FEB"),),
    ),
    ColorMatchCase(
        name="currency-green-anchor",
        source_segments=(
            SourceSegment("Total: ", "111111"),
            SourceSegment("$1,500", "1F8A3B"),
        ),
        translation="합계: $1,500",
        anchors=(ColorAnchor("$1,500", "1F8A3B"),),
    ),
    ColorMatchCase(
        name="emphasis-red-reordered",
        source_segments=(
            SourceSegment("strongly ", "D72638"),
            SourceSegment("recommend this approach", "111111"),
        ),
        translation="이 접근 방식을 강력히 추천합니다",
        anchors=(ColorAnchor("강력히", "D72638"),),
    ),
)


def normalize_color(color: str | None) -> str | None:
    """Normalize a color value to uppercase hex without a leading '#'. """
    if color is None:
        return None
    return str(color).replace("#", "").upper()


def build_color_match_fixture(
    path: Path,
    cases: Sequence[ColorMatchCase] = DEFAULT_COLOR_MATCH_CASES,
) -> None:
    """Create a PPTX containing source paragraphs with colored runs."""
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    text_box = slide.shapes.add_textbox(
        Inches(0.7),
        Inches(0.7),
        Inches(8.6),
        Inches(4.6),
    )
    text_frame = text_box.text_frame
    text_frame.clear()

    for case_index, case in enumerate(cases):
        paragraph = text_frame.paragraphs[0] if case_index == 0 else text_frame.add_paragraph()
        paragraph.text = ""
        for segment in case.source_segments:
            run = paragraph.add_run()
            run.text = segment.text
            run.font.size = Pt(20)
            run.font.color.rgb = RGBColor.from_string(normalize_color(segment.color) or "111111")

    presentation.save(path)


def collect_paragraph_run_segments(path: Path) -> list[list[tuple[str, str | None]]]:
    """Return visible text/color segments for each non-empty paragraph."""
    presentation = Presentation(path)
    paragraphs: list[list[tuple[str, str | None]]] = []

    for slide in presentation.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                segments: list[tuple[str, str | None]] = []
                for run in paragraph.runs:
                    if not run.text:
                        continue
                    color = None
                    try:
                        color = normalize_color(run.font.color.rgb)
                    except AttributeError:
                        color = None
                    segments.append((run.text, color))
                if segments:
                    paragraphs.append(segments)

    return paragraphs


def verify_output_segments(
    cases: Sequence[ColorMatchCase],
    output_segments: Sequence[Sequence[tuple[str, str | None]]],
) -> list[ColorMatchCaseResult]:
    """Verify output text and anchor colors for all cases."""
    results: list[ColorMatchCaseResult] = []

    for case_index, case in enumerate(cases):
        if case_index >= len(output_segments):
            results.append(
                ColorMatchCaseResult(
                    name=case.name,
                    passed=False,
                    output_text="",
                    errors=["missing output paragraph"],
                )
            )
            continue

        segments = output_segments[case_index]
        result = verify_case_segments(case, segments)
        results.append(result)

    return results


def verify_case_segments(
    case: ColorMatchCase,
    segments: Sequence[tuple[str, str | None]],
) -> ColorMatchCaseResult:
    """Verify one case against output run text/color segments."""
    output_text = "".join(text for text, _ in segments)
    errors: list[str] = []

    if output_text != case.translation:
        errors.append(
            f"text mismatch: expected {case.translation!r}, got {output_text!r}"
        )

    char_colors: list[str | None] = []
    for text, color in segments:
        char_colors.extend([normalize_color(color)] * len(text))

    for anchor in case.anchors:
        expected_color = normalize_color(anchor.color)
        start = output_text.find(anchor.text)
        if start < 0:
            errors.append(f"anchor {anchor.text!r} not found")
            continue
        end = start + len(anchor.text)
        actual_colors = {
            color
            for char, color in zip(output_text[start:end], char_colors[start:end])
            if not char.isspace()
        }
        if actual_colors != {expected_color}:
            errors.append(
                f"anchor {anchor.text!r} color mismatch: expected {expected_color}, got {sorted(actual_colors)}"
            )

    return ColorMatchCaseResult(
        name=case.name,
        passed=not errors,
        output_text=output_text,
        errors=errors,
    )


async def verify_color_matching_once(
    provider: str,
    model: str | None = None,
    *,
    iteration: int = 1,
    cases: Sequence[ColorMatchCase] = DEFAULT_COLOR_MATCH_CASES,
    output_dir: Path | None = None,
) -> ColorMatchIterationResult:
    """Run one end-to-end color matching verification iteration."""
    model_name = model or DEFAULT_LIGHT_MODEL[provider]

    if output_dir is None:
        with tempfile.TemporaryDirectory(prefix=f"color-match-{provider}-") as tmp:
            result = await _verify_color_matching_once_in_dir(
                provider=provider,
                model=model_name,
                iteration=iteration,
                cases=cases,
                output_dir=Path(tmp),
            )
            result.input_path = None
            result.output_path = None
            return result

    output_dir.mkdir(parents=True, exist_ok=True)
    return await _verify_color_matching_once_in_dir(
        provider=provider,
        model=model_name,
        iteration=iteration,
        cases=cases,
        output_dir=output_dir,
    )


async def _verify_color_matching_once_in_dir(
    *,
    provider: str,
    model: str,
    iteration: int,
    cases: Sequence[ColorMatchCase],
    output_dir: Path,
) -> ColorMatchIterationResult:
    input_path = output_dir / f"{provider}-color-input-{iteration}.pptx"
    output_path = output_dir / f"{provider}-color-output-{iteration}.pptx"

    build_color_match_fixture(input_path, cases)
    with input_path.open("rb") as input_file:
        paragraphs, presentation = PPTParser().extract_paragraphs(input_file)

    translations = [case.translation for case in cases]
    distributions = await TranslationService._fix_color_distributions_async(
        paragraphs,
        translations,
        provider,
        model_name=model,
    )

    PPTWriter().apply_translations(
        paragraphs,
        translations,
        presentation,
        text_fit_mode="none",
        color_distributions=distributions,
        output_path=output_path,
    )

    output_segments = collect_paragraph_run_segments(output_path)
    case_results = verify_output_segments(cases, output_segments)

    if distributions is None:
        for result in case_results:
            result.passed = False
            result.errors.append("no validated color distribution returned")
    else:
        for index, result in enumerate(case_results):
            if index not in distributions:
                result.passed = False
                result.errors.append("missing validated color distribution for paragraph")

    return ColorMatchIterationResult(
        provider=provider,
        model=model,
        iteration=iteration,
        passed=all(result.passed for result in case_results),
        cases=case_results,
        input_path=input_path,
        output_path=output_path,
    )


async def verify_color_matching_loop(
    providers: Iterable[str],
    *,
    iterations: int = 1,
    models: dict[str, str | None] | None = None,
    output_dir: Path | None = None,
) -> list[ColorMatchIterationResult]:
    """Run color matching verification for providers over N iterations."""
    results: list[ColorMatchIterationResult] = []
    models = models or {}

    for provider in providers:
        for iteration in range(1, max(1, iterations) + 1):
            provider_output_dir = None
            if output_dir is not None:
                provider_output_dir = output_dir / provider
            result = await verify_color_matching_once(
                provider,
                models.get(provider),
                iteration=iteration,
                output_dir=provider_output_dir,
            )
            results.append(result)

    return results


def run_color_matching_loop(
    providers: Iterable[str],
    *,
    iterations: int = 1,
    models: dict[str, str | None] | None = None,
    output_dir: Path | None = None,
) -> list[ColorMatchIterationResult]:
    """Synchronous wrapper for the async verification loop."""
    return asyncio.run(
        verify_color_matching_loop(
            providers,
            iterations=iterations,
            models=models,
            output_dir=output_dir,
        )
    )
