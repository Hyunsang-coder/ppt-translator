"""Tests for the FastAPI endpoints."""

from __future__ import annotations

import io
import zipfile
from unittest.mock import MagicMock, patch

import pytest
from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.util import Inches

from api import (
    app,
    FilenameSettings,
    SUPPORTED_LANGUAGES,
    SUPPORTED_MODELS,
    generate_output_filename,
)
from src.core.ppt_parser import PPTParser
from src.services.job_manager import JobState, JobType, get_job_manager
from src.services.review_session import ReviewSession
from src.utils.config import get_settings


@pytest.fixture
def client():
    """Create a test client."""
    return TestClient(app)


@pytest.fixture
def sample_pptx_bytes():
    """Create a minimal valid PPTX file in memory."""
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as zf:
        # Minimal PPTX structure
        zf.writestr("[Content_Types].xml", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
    <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
    <Default Extension="xml" ContentType="application/xml"/>
    <Override PartName="/ppt/presentation.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/>
</Types>""")
        zf.writestr("_rels/.rels", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
    <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="ppt/presentation.xml"/>
</Relationships>""")
        zf.writestr("ppt/presentation.xml", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:presentation xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
    <p:sldIdLst/>
</p:presentation>""")
        zf.writestr("ppt/_rels/presentation.xml.rels", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
</Relationships>""")
    buffer.seek(0)
    return buffer.read()


@pytest.fixture
def review_job():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    box.text_frame.text = "원문"
    source = io.BytesIO()
    presentation.save(source)
    source_bytes = source.getvalue()
    paragraphs, parsed = PPTParser().extract_paragraphs(io.BytesIO(source_bytes))
    session = ReviewSession(
        presentation=parsed,
        paragraphs=paragraphs,
        translated_texts=["OLD"],
        findings=[],
        source_lang="한국어",
        target_lang="영어",
        model="stub",
        source_pptx=source_bytes,
    )
    manager = get_job_manager()
    job = manager.create_job(JobType.TRANSLATION)
    job.state = JobState.COMPLETED
    job.review_session = session
    job.output_file = io.BytesIO(b"published-before-review")
    job.output_filename = "review.pptx"
    yield job
    manager._jobs.pop(job.id, None)


class TestFilenameGeneration:
    """Tests for output filename generation."""

    def test_generate_output_filename_respects_component_order(self):
        """Custom component order should be reflected in output."""
        settings = FilenameSettings(
            mode="auto",
            includeLanguage=True,
            includeOriginalName=True,
            includeModel=False,
            includeDate=False,
            componentOrder=["originalName", "language", "model", "date"],
        )

        filename = generate_output_filename(
            filename_settings=settings,
            original_filename="original.pptx",
            target_language="영어",
            model="gpt-5.6-sol",
        )

        assert filename == "original_EN.pptx"

    def test_generate_output_filename_uses_default_order_when_missing(self):
        """Older clients without componentOrder should keep default behavior."""
        settings = FilenameSettings(
            mode="auto",
            includeLanguage=True,
            includeOriginalName=True,
            includeModel=False,
            includeDate=False,
        )

        filename = generate_output_filename(
            filename_settings=settings,
            original_filename="original.pptx",
            target_language="영어",
            model="gpt-5.6-sol",
        )

        assert filename == "EN_original.pptx"


class TestHealthEndpoint:
    """Tests for the health check endpoint."""

    def test_health_check(self, client):
        """Test health check returns expected fields."""
        response = client.get("/health")
        assert response.status_code == 200
        data = response.json()
        assert data["status"] == "healthy"
        assert "timestamp" in data
        assert "openai_api_key_configured" in data
        assert "anthropic_api_key_configured" in data

    def test_health_check_includes_jobs_info(self, client):
        """Test health check includes job concurrency info."""
        response = client.get("/health")
        assert response.status_code == 200
        data = response.json()
        assert "jobs" in data
        jobs = data["jobs"]
        assert "running" in jobs
        assert "pending" in jobs
        assert "max_running" in jobs
        assert isinstance(jobs["running"], int)
        assert isinstance(jobs["max_running"], int)

    def test_health_check_includes_memory(self, client):
        """Test health check includes memory usage."""
        response = client.get("/health")
        assert response.status_code == 200
        data = response.json()
        assert "memory_usage_mb" in data
        assert isinstance(data["memory_usage_mb"], (int, float))


class TestConfigEndpoints:
    """Tests for configuration endpoints."""

    def test_get_models_all(self, client):
        """Test getting all models."""
        response = client.get("/api/v1/models")
        assert response.status_code == 200
        data = response.json()
        assert "models" in data
        assert len(data["models"]) > 0
        # Check model structure
        model = data["models"][0]
        assert "id" in model
        assert "name" in model
        assert "provider" in model

    def test_get_models_by_provider(self, client):
        """Test getting models filtered by provider."""
        response = client.get("/api/v1/models?provider=openai")
        assert response.status_code == 200
        data = response.json()
        assert all(m["provider"] == "openai" for m in data["models"])

        response = client.get("/api/v1/models?provider=anthropic")
        assert response.status_code == 200
        data = response.json()
        assert all(m["provider"] == "anthropic" for m in data["models"])

    def test_get_models_invalid_provider(self, client):
        """Test getting models with invalid provider."""
        response = client.get("/api/v1/models?provider=invalid")
        assert response.status_code == 400

    def test_get_languages(self, client):
        """Test getting supported languages."""
        response = client.get("/api/v1/languages")
        assert response.status_code == 200
        data = response.json()
        assert "languages" in data
        assert len(data["languages"]) > 0
        # Check language structure
        lang = data["languages"][0]
        assert "code" in lang
        assert "name" in lang
        # Check Auto is present
        codes = [l["code"] for l in data["languages"]]
        assert "Auto" in codes

    def test_get_config(self, client):
        """Test getting application config."""
        response = client.get("/api/v1/config")
        assert response.status_code == 200
        data = response.json()
        assert "max_upload_size_mb" in data
        assert "providers" in data
        assert "default_provider" in data
        assert "default_model" in data


class TestJobEndpoints:
    """Tests for job management endpoints."""

    def test_create_job_invalid_file_type(self, client):
        """Test creating job with invalid file type."""
        response = client.post(
            "/api/v1/jobs",
            files={"ppt_file": ("test.txt", b"not a pptx", "text/plain")},
            data={"provider": "openai", "model": "gpt-5.6-sol"},
        )
        assert response.status_code == 400
        assert "Invalid file type" in response.json()["detail"]

    def test_create_job_validation_failure_releases_reserved_slot(self, client):
        """Invalid uploads should not consume active job capacity."""
        from src.services.job_manager import get_job_manager

        job_manager = get_job_manager()
        active_before = job_manager.get_active_count()

        response = client.post(
            "/api/v1/jobs",
            files={"ppt_file": ("test.pptx", b"not a real pptx", "application/octet-stream")},
            data={"provider": "openai", "model": "gpt-5.6-sol"},
        )

        assert response.status_code == 400
        assert job_manager.get_active_count() == active_before

    def test_create_job_invalid_provider(self, client, sample_pptx_bytes):
        """Test creating job with invalid provider."""
        response = client.post(
            "/api/v1/jobs",
            files={"ppt_file": ("test.pptx", sample_pptx_bytes, "application/octet-stream")},
            data={"provider": "invalid", "model": "test"},
        )
        assert response.status_code == 400
        assert "Invalid provider" in response.json()["detail"]

    def test_create_job_invalid_model_for_provider(self, client, sample_pptx_bytes):
        """A model not in the provider's allowlist must be rejected (no key abuse)."""
        response = client.post(
            "/api/v1/jobs",
            files={"ppt_file": ("test.pptx", sample_pptx_bytes, "application/octet-stream")},
            data={"provider": "anthropic", "model": "claude-3-opus-20240229"},
        )
        assert response.status_code == 400
        assert "Invalid model" in response.json()["detail"]

    def test_create_job_invalid_glossary_json(self, client, sample_pptx_bytes):
        """Malformed glossary_json must 400 before the job starts translating."""
        response = client.post(
            "/api/v1/jobs",
            files={"ppt_file": ("test.pptx", sample_pptx_bytes, "application/octet-stream")},
            data={
                "provider": "openai",
                "model": "gpt-5.6-sol",
                "glossary_json": "{not-json",
            },
        )
        assert response.status_code == 400
        assert "Glossary error" in response.json()["detail"]

    def test_parse_glossary_excel(self, client):
        """Excel import endpoint returns structured entries and skips header row."""
        import pandas as pd

        buffer = io.BytesIO()
        pd.DataFrame(
            [["원문", "번역"], ["PUBG", "배틀그라운드"], ["Stream", "스트리밍"]]
        ).to_excel(buffer, index=False, header=False)
        buffer.seek(0)

        response = client.post(
            "/api/v1/glossary/parse",
            files={
                "glossary_file": (
                    "glossary.xlsx",
                    buffer.getvalue(),
                    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                )
            },
        )
        assert response.status_code == 200
        data = response.json()
        assert data["count"] == 2
        assert {"source": "PUBG", "target": "배틀그라운드"} in data["entries"]
        assert all(e["source"] != "원문" for e in data["entries"])

    def test_get_job_not_found(self, client):
        """Test getting non-existent job."""
        response = client.get("/api/v1/jobs/non-existent-id")
        assert response.status_code == 404

    def test_delete_job_not_found(self, client):
        """Test deleting non-existent job."""
        response = client.delete("/api/v1/jobs/non-existent-id")
        assert response.status_code == 404

    def test_job_result_not_found(self, client):
        """Test downloading result of non-existent job."""
        response = client.get("/api/v1/jobs/non-existent-id/result")
        assert response.status_code == 404

    def test_create_job_returns_429_when_queue_full(self, client, sample_pptx_bytes):
        """Test that creating a job returns 429 when active jobs exceed limit."""
        from src.services.job_manager import get_job_manager, JobState

        job_manager = get_job_manager()
        settings = get_settings()
        max_allowed = settings.max_running_jobs + settings.max_queued_jobs

        # Fill up with fake pending jobs
        fake_jobs = []
        for _ in range(max_allowed):
            job = job_manager.create_job(JobType.TRANSLATION)
            fake_jobs.append(job)

        try:
            response = client.post(
                "/api/v1/jobs",
                files={"ppt_file": ("test.pptx", sample_pptx_bytes, "application/octet-stream")},
                data={"provider": "openai", "model": "gpt-5.6-sol"},
            )
            assert response.status_code == 429
            assert "바쁩니다" in response.json()["detail"]
        finally:
            # Clean up fake jobs so they don't affect other tests
            for job in fake_jobs:
                job.state = JobState.COMPLETED
                job.completed_at = 0  # mark as old
            job_manager._cleanup_old_jobs()


class TestExtractionEndpoint:
    """Tests for text extraction endpoint."""

    def test_extract_invalid_file_type(self, client):
        """Test extraction with invalid file type."""
        response = client.post(
            "/api/v1/extract",
            files={"ppt_file": ("test.txt", b"not a pptx", "text/plain")},
        )
        assert response.status_code == 400
        assert "Invalid file type" in response.json()["detail"]

    def test_extract_valid_pptx(self, client, sample_pptx_bytes):
        """Test extraction with valid PPTX."""
        response = client.post(
            "/api/v1/extract",
            files={"ppt_file": ("test.pptx", sample_pptx_bytes, "application/octet-stream")},
            data={"figures": "omit", "charts": "labels", "with_notes": "false"},
        )
        assert response.status_code == 200
        data = response.json()
        assert "markdown" in data
        assert "slide_count" in data


class TestReviewEndpoints:
    def test_proposal_is_previewed_then_staged_then_committed(self, client, review_job):
        response = client.get(f"/api/v1/jobs/{review_job.id}/fragments")
        assert response.status_code == 200
        assert response.json()["revision"] == 0
        assert response.json()["dirty"] is False

        response = client.post(
            f"/api/v1/jobs/{review_job.id}/fragments/0/proposals",
            json={"action": "edit", "target": "NEW"},
        )
        assert response.status_code == 200
        proposal = response.json()
        assert proposal["old_target"] == "OLD"
        assert proposal["target"] == "NEW"
        assert review_job.review_session.translated_texts[0] == "OLD"

        with patch("api._record_edit"):
            response = client.post(
                f"/api/v1/jobs/{review_job.id}/proposals/{proposal['proposal_id']}/apply",
                json={"expected_revision": 0},
            )
        assert response.status_code == 200
        assert response.json()["revision"] == 1
        assert review_job.review_session.translated_texts[0] == "NEW"
        # Staging never changes the published download.
        assert review_job.output_file.getvalue() == b"published-before-review"

        response = client.post(
            f"/api/v1/jobs/{review_job.id}/review/commit",
            json={"expected_revision": 1},
        )
        assert response.status_code == 200
        assert response.json()["dirty"] is False
        rendered = Presentation(review_job.output_file)
        assert rendered.slides[0].shapes[0].text == "NEW"

    def test_stale_apply_returns_conflict(self, client, review_job):
        response = client.post(
            f"/api/v1/jobs/{review_job.id}/fragments/0/proposals",
            json={"action": "edit", "target": "NEW"},
        )
        proposal_id = response.json()["proposal_id"]
        review_job.review_session.apply_edit(0, "OTHER")
        response = client.post(
            f"/api/v1/jobs/{review_job.id}/proposals/{proposal_id}/apply",
            json={"expected_revision": 0},
        )
        assert response.status_code == 409

    def test_failed_commit_keeps_published_file(self, client, review_job):
        review_job.review_session.apply_edit(0, "NEW")
        before = review_job.output_file.getvalue()
        with patch.object(review_job.review_session, "render", side_effect=RuntimeError("boom")):
            response = client.post(
                f"/api/v1/jobs/{review_job.id}/review/commit",
                json={"expected_revision": 1},
            )
        assert response.status_code == 500
        assert review_job.output_file.getvalue() == before
        assert review_job.review_session.dirty is True


class TestJobManager:
    """Tests for job manager functionality."""

    def test_job_manager_singleton(self):
        """Test job manager returns same instance."""
        from src.services import get_job_manager

        manager1 = get_job_manager()
        manager2 = get_job_manager()
        assert manager1 is manager2

    def test_job_creation(self):
        """Test job creation and retrieval."""
        from src.services import get_job_manager, JobType

        manager = get_job_manager()
        job = manager.create_job(JobType.TRANSLATION)

        assert job.id is not None
        assert job.job_type == JobType.TRANSLATION

        retrieved = manager.get_job(job.id)
        assert retrieved is not None
        assert retrieved.id == job.id

    @pytest.mark.asyncio
    async def test_job_deletion(self):
        """Test job cancellation (job stays in store with cancelled state)."""
        from src.services import get_job_manager, JobType, JobState

        manager = get_job_manager()
        job = manager.create_job(JobType.TRANSLATION)
        job_id = job.id

        assert await manager.delete_job(job_id) is True
        # Job remains in store with cancelled state (cleaned up later)
        cancelled_job = manager.get_job(job_id)
        assert cancelled_job is not None
        assert cancelled_job.state == JobState.CANCELLED
        assert cancelled_job.completed_at is not None
