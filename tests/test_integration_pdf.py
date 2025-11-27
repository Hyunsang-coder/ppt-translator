"""
Integration tests for PDF to PPT conversion using real PDF files.

These tests require:
- OPENAI_API_KEY environment variable set (for Vision API tests only)
- test-assets/ directory with PDF files

Run with: pytest tests/test_integration_pdf.py -v -s
"""

import io
import os
import unittest
from pathlib import Path

import pytest

try:
    from PIL import Image
    from pptx import Presentation
    from pptx.util import Emu
except ImportError as e:
    pytest.skip(f"Required package not installed: {e}", allow_module_level=True)

from src.core.pdf_processor import PDFProcessor, PageOCRResult, TextBlock
from src.core.pdf_to_ppt_writer import PDFToPPTWriter, TextBoxStyle


# Helper to check if API key is available
def requires_api_key(test_func):
    """Decorator to skip tests that require OPENAI_API_KEY."""
    return pytest.mark.skipif(
        not os.environ.get("OPENAI_API_KEY"),
        reason="OPENAI_API_KEY not set"
    )(test_func)


class TestPDFBasicOperations(unittest.TestCase):
    """Basic PDF operations that don't require API key."""

    @classmethod
    def setUpClass(cls):
        """Set up test fixtures."""
        cls.test_assets_dir = Path(__file__).parent.parent / "test-assets"
        cls.pdf_files = list(cls.test_assets_dir.glob("*.pdf"))

    def test_pdf_files_exist(self):
        """Verify test PDF files are available."""
        self.assertTrue(
            len(self.pdf_files) > 0,
            f"No PDF files found in {self.test_assets_dir}"
        )
        print(f"\nFound {len(self.pdf_files)} PDF files:")
        for pdf in self.pdf_files:
            print(f"  - {pdf.name} ({pdf.stat().st_size / 1024:.1f} KB)")

    def test_pdf_to_images_conversion(self):
        """Test PDF to image conversion without Vision API."""
        if not self.pdf_files:
            self.skipTest("No test PDFs available")
            
        # Use smaller PDF for faster testing
        pdf_path = min(self.pdf_files, key=lambda p: p.stat().st_size)
        
        # Create processor without API key (only for image conversion)
        processor = PDFProcessor(api_key="", model="gpt-4o", dpi=150)
        
        with open(pdf_path, "rb") as f:
            pdf_buffer = io.BytesIO(f.read())
        
        images = processor.convert_pdf_to_images(pdf_buffer, max_pages=1)
        
        self.assertEqual(len(images), 1)
        page_num, image = images[0]
        self.assertEqual(page_num, 1)
        self.assertIsInstance(image, Image.Image)
        
        print(f"\nConverted {pdf_path.name}:")
        print(f"  - Page 1: {image.width}x{image.height} pixels")

    def test_ppt_creation_with_mock_blocks(self):
        """Test PPT creation with manually created TextBlocks (no API needed)."""
        if not self.pdf_files:
            self.skipTest("No test PDFs available")
            
        pdf_path = min(self.pdf_files, key=lambda p: p.stat().st_size)
        processor = PDFProcessor(api_key="", model="gpt-4o", dpi=150)
        
        with open(pdf_path, "rb") as f:
            pdf_buffer = io.BytesIO(f.read())
        
        images = processor.convert_pdf_to_images(pdf_buffer, max_pages=1)
        page_num, image = images[0]
        
        # Create mock text blocks for testing
        mock_blocks = [
            TextBlock(
                text="Title Text",
                left=100, top=50, width=400, height=60,
                block_type="title"
            ),
            TextBlock(
                text="Body text content here. This is a paragraph.",
                left=100, top=150, width=400, height=100,
                block_type="body"
            ),
            TextBlock(
                text="한글 텍스트 테스트",
                left=100, top=280, width=300, height=50,
                block_type="body"
            ),
        ]
        
        ocr_result = PageOCRResult(
            page_number=page_num,
            image=image,
            text_blocks=mock_blocks,
        )
        
        writer = PDFToPPTWriter()
        ppt_buffer = writer.create_presentation([ocr_result], include_background=True)
        
        # Verify PPT was created
        prs = Presentation(ppt_buffer)
        self.assertEqual(len(prs.slides), 1)
        
        # Should have background image + 3 text boxes
        shapes = list(prs.slides[0].shapes)
        self.assertEqual(len(shapes), 4)
        
        print(f"\n✅ PPT created with {len(shapes)} shapes (1 background + 3 text boxes)")


@pytest.mark.skipif(
    not os.environ.get("OPENAI_API_KEY"),
    reason="OPENAI_API_KEY not set"
)
class TestIntegrationPDFToPPT(unittest.TestCase):
    """Integration tests using real PDF files from test-assets/ (requires API key)"""

    @classmethod
    def setUpClass(cls):
        """Set up test fixtures."""
        cls.test_assets_dir = Path(__file__).parent.parent / "test-assets"
        cls.output_dir = Path(__file__).parent.parent / "test-output"
        cls.output_dir.mkdir(exist_ok=True)
        
        cls.api_key = os.environ.get("OPENAI_API_KEY", "")
        cls.pdf_files = list(cls.test_assets_dir.glob("*.pdf"))
        
    def setUp(self):
        """Set up for each test."""
        self.processor = PDFProcessor(
            api_key=self.api_key,
            model="gpt-4o",
            dpi=150
        )
        self.writer = PDFToPPTWriter()

    @pytest.mark.slow
    def test_full_pipeline_pdf_test_1(self):
        """Full integration test with pdf_test_1.pdf (Vision API call)."""
        pdf_path = self.test_assets_dir / "pdf_test_1.pdf"
        if not pdf_path.exists():
            self.skipTest("pdf_test_1.pdf not found")
            
        self._run_full_pipeline(pdf_path, max_pages=1)

    @pytest.mark.slow  
    def test_full_pipeline_pdf_test_2(self):
        """Full integration test with pdf_test_2.pdf (Vision API call)."""
        pdf_path = self.test_assets_dir / "pdf_test_2.pdf"
        if not pdf_path.exists():
            self.skipTest("pdf_test_2.pdf not found")
            
        self._run_full_pipeline(pdf_path, max_pages=1)

    def _run_full_pipeline(self, pdf_path: Path, max_pages: int = 1):
        """Run full PDF -> Vision -> PPT pipeline and save output."""
        print(f"\n{'='*60}")
        print(f"Processing: {pdf_path.name}")
        print(f"{'='*60}")
        
        # Step 1: Load PDF
        with open(pdf_path, "rb") as f:
            pdf_buffer = io.BytesIO(f.read())
        
        # Step 2: Process with Vision
        ocr_results = self.processor.process_pdf(pdf_buffer, max_pages=max_pages)
        
        self.assertEqual(len(ocr_results), max_pages)
        
        # Step 3: Analyze detected blocks
        for result in ocr_results:
            print(f"\nPage {result.page_number}:")
            print(f"  Image: {result.image_width}x{result.image_height}")
            print(f"  Blocks detected: {len(result.text_blocks)}")
            
            # Check for potential overlaps
            overlaps = self._find_overlapping_blocks(result.text_blocks)
            if overlaps:
                print(f"  ⚠️  Overlapping blocks: {len(overlaps)}")
                for (i, j), overlap_area in overlaps:
                    b1, b2 = result.text_blocks[i], result.text_blocks[j]
                    print(f"      Block {i} ({b1.block_type}) vs Block {j} ({b2.block_type})")
                    print(f"      Overlap area: {overlap_area} px²")
            
            for idx, block in enumerate(result.text_blocks):
                text_preview = block.text[:40] + "..." if len(block.text) > 40 else block.text
                print(f"    [{idx}] {block.block_type}: '{text_preview}'")
                print(f"        Position: ({block.left}, {block.top}) Size: {block.width}x{block.height}")
        
        # Step 4: Create PPT
        ppt_buffer = self.writer.create_presentation(ocr_results, include_background=True)
        
        # Step 5: Verify and save PPT
        prs = Presentation(ppt_buffer)
        self.assertEqual(len(prs.slides), max_pages)
        
        output_path = self.output_dir / f"{pdf_path.stem}_output.pptx"
        with open(output_path, "wb") as f:
            ppt_buffer.seek(0)
            f.write(ppt_buffer.read())
        
        print(f"\n✅ Output saved: {output_path}")
        print(f"   Slide size: {prs.slide_width} x {prs.slide_height} EMU")
        
        # Verify shapes
        for slide_idx, slide in enumerate(prs.slides):
            shapes = list(slide.shapes)
            print(f"   Slide {slide_idx + 1}: {len(shapes)} shapes")

    def _find_overlapping_blocks(self, blocks: list) -> list:
        """Find pairs of overlapping blocks and their overlap areas."""
        overlaps = []
        for i in range(len(blocks)):
            for j in range(i + 1, len(blocks)):
                b1, b2 = blocks[i], blocks[j]
                overlap_area = self._calculate_overlap_area(b1, b2)
                if overlap_area > 0:
                    overlaps.append(((i, j), overlap_area))
        return overlaps

    def _calculate_overlap_area(self, b1: TextBlock, b2: TextBlock) -> int:
        """Calculate overlap area between two blocks in pixels."""
        x_overlap = max(0, min(b1.right, b2.right) - max(b1.left, b2.left))
        y_overlap = max(0, min(b1.bottom, b2.bottom) - max(b1.top, b2.top))
        return x_overlap * y_overlap


class TestTextBlockOverlapDetection(unittest.TestCase):
    """Unit tests for overlap detection logic (no API calls)."""

    def test_no_overlap_horizontal(self):
        """Test blocks side by side don't overlap."""
        b1 = TextBlock("A", left=0, top=0, width=100, height=50)
        b2 = TextBlock("B", left=110, top=0, width=100, height=50)
        
        overlap = self._calculate_overlap_area(b1, b2)
        self.assertEqual(overlap, 0)

    def test_no_overlap_vertical(self):
        """Test blocks stacked vertically don't overlap."""
        b1 = TextBlock("A", left=0, top=0, width=100, height=50)
        b2 = TextBlock("B", left=0, top=60, width=100, height=50)
        
        overlap = self._calculate_overlap_area(b1, b2)
        self.assertEqual(overlap, 0)

    def test_partial_overlap(self):
        """Test partially overlapping blocks."""
        b1 = TextBlock("A", left=0, top=0, width=100, height=100)
        b2 = TextBlock("B", left=50, top=50, width=100, height=100)
        
        # Overlap: 50x50 = 2500
        overlap = self._calculate_overlap_area(b1, b2)
        self.assertEqual(overlap, 2500)

    def test_complete_containment(self):
        """Test when one block completely contains another."""
        b1 = TextBlock("A", left=0, top=0, width=200, height=200)
        b2 = TextBlock("B", left=50, top=50, width=50, height=50)
        
        # b2 is completely inside b1: 50x50 = 2500
        overlap = self._calculate_overlap_area(b1, b2)
        self.assertEqual(overlap, 2500)

    def test_adjacent_blocks_no_overlap(self):
        """Test blocks that touch but don't overlap."""
        b1 = TextBlock("A", left=0, top=0, width=100, height=100)
        b2 = TextBlock("B", left=100, top=0, width=100, height=100)
        
        overlap = self._calculate_overlap_area(b1, b2)
        self.assertEqual(overlap, 0)

    def _calculate_overlap_area(self, b1: TextBlock, b2: TextBlock) -> int:
        """Calculate overlap area between two blocks."""
        x_overlap = max(0, min(b1.right, b2.right) - max(b1.left, b2.left))
        y_overlap = max(0, min(b1.bottom, b2.bottom) - max(b1.top, b2.top))
        return x_overlap * y_overlap


if __name__ == "__main__":
    unittest.main()

