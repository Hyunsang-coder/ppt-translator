
import io
import unittest
from unittest.mock import MagicMock, patch
try:
    from PIL import Image
except ImportError:
    Image = None

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor

from src.core.pdf_processor import PageOCRResult, TextBlock
from src.core.pdf_to_ppt_writer import PDFToPPTWriter, TextBoxStyle, LayoutRect

class TestPDFToPPTLayout(unittest.TestCase):
    def setUp(self):
        self.writer = PDFToPPTWriter()

    def test_portrait_mode_slide_size(self):
        """Test that a portrait input image results in a portrait slide."""
        if not Image:
            self.skipTest("PIL not installed")

        width, height = 100, 200
        image = Image.new("RGB", (width, height), color="white")
        
        ocr_result = PageOCRResult(
            page_number=1,
            image=image,
            text_blocks=[],
            image_width=width,
            image_height=height
        )

        ppt_buffer = self.writer.create_presentation([ocr_result])
        prs = Presentation(ppt_buffer)
        
        self.assertTrue(prs.slide_height > prs.slide_width, 
                        f"Slide should be portrait. Width: {prs.slide_width}, Height: {prs.slide_height}")

    def test_mixed_orientation_layout(self):
        """Test that mixed orientation pages are handled without crashing."""
        if not Image:
            self.skipTest("PIL not installed")

        img1 = Image.new("RGB", (100, 200), "white")
        res1 = PageOCRResult(1, img1, [], 100, 200)

        img2 = Image.new("RGB", (200, 100), "white")
        res2 = PageOCRResult(2, img2, [], 200, 100)

        ppt_buffer = self.writer.create_presentation([res1, res2])
        prs = Presentation(ppt_buffer)

        self.assertTrue(prs.slide_height > prs.slide_width)
        self.assertEqual(len(prs.slides), 2)

    def test_color_enforcement(self):
        """Test that user-defined colors override Vision-detected colors."""
        if not Image:
            self.skipTest("PIL not installed")

        # User wants White bg, Black text (Explicitly set, not None)
        style = TextBoxStyle(
            background_color=(255, 255, 255),
            text_color=(0, 0, 0)
        )
        writer = PDFToPPTWriter(text_style=style)

        # Vision detected Red text
        block = TextBlock(
            text="Hello",
            left=0, top=0, width=100, height=50,
            text_color=(255, 0, 0) # Red
        )
        
        width, height = 200, 200
        image = Image.new("RGB", (width, height), "white")
        res = PageOCRResult(1, image, [block], width, height)

        # Test 1: White Image -> Expect White BG, Black Text
        ppt_buffer = writer.create_presentation([res])
        prs = Presentation(ppt_buffer)
        shape = prs.slides[0].shapes[-1]
        
        run = shape.text_frame.paragraphs[0].runs[0]
        self.assertEqual(run.font.color.rgb, RGBColor(0, 0, 0)) # Black text on White
        self.assertEqual(shape.fill.fore_color.rgb, RGBColor(255, 255, 255)) # White BG
        
        # Test 2: Black Image -> Expect White BG, Black Text (BECAUSE STYLE IS ENFORCED)
        # Even though image is black, we forced White BG in style.
        image_black = Image.new("RGB", (width, height), "black")
        res_black = PageOCRResult(1, image_black, [block], width, height)
        
        ppt_buffer_black = writer.create_presentation([res_black])
        prs_black = Presentation(ppt_buffer_black)
        shape_black = prs_black.slides[0].shapes[-1]
        
        run_black = shape_black.text_frame.paragraphs[0].runs[0]
        
        # This was failing because it was adapting to Black image (White text, Black BG)
        # But now we fixed logic to respect style if set.
        # So it should be Black Text (0,0,0) on White BG (255,255,255) as per style.
        self.assertEqual(run_black.font.color.rgb, RGBColor(0, 0, 0))
        self.assertEqual(shape_black.fill.fore_color.rgb, RGBColor(255, 255, 255))

    def test_text_box_preserves_original_width(self):
        """Test that text box width is preserved (and expanded) to fully cover original text region."""
        if not Image:
            self.skipTest("PIL not installed")
            
        writer = PDFToPPTWriter()
        
        # Short text ("Title"), but Vision gave a wide box (1000px)
        # The box should NOT be shrunk - we need to cover the original text
        wide_width = 1000
        block = TextBlock(
            text="Title",
            left=0, top=0, width=wide_width, height=50,
            block_type="title"
        )
        
        width, height = 2000, 2000
        image = Image.new("RGB", (width, height), "white")
        res = PageOCRResult(1, image, [block], width, height)

        ppt_buffer = writer.create_presentation([res])
        prs = Presentation(ppt_buffer)
        shape = prs.slides[0].shapes[-1]
        
        # Calculate expected width in EMU
        # block.width (1000) * scale (1.0) * 9525 EMU/px
        original_width_emu = 1000 * 9525
        # Expansion: 5% left + 5% right = 10%
        expected_width = int(original_width_emu * 1.1)
        
        # Allow small rounding differences
        self.assertAlmostEqual(shape.width, expected_width, delta=100,
                        msg="Text box width should be expanded 10% to cover original text region")

    def test_text_box_top_left_positioning(self):
        """Test that text box is positioned relative to Top-Left of original region (with expansion)."""
        if not Image:
            self.skipTest("PIL not installed")
            
        writer = PDFToPPTWriter()
        
        # Text block at specific position
        block = TextBlock(
            text="Centered Title",
            left=200, top=100, width=600, height=80,
            block_type="title"  # Center aligned
        )
        
        width, height = 1000, 1000
        image = Image.new("RGB", (width, height), "white")
        res = PageOCRResult(1, image, [block], width, height)

        ppt_buffer = writer.create_presentation([res])
        prs = Presentation(ppt_buffer)
        shape = prs.slides[0].shapes[-1]
        
        # Expected position in EMU (Top-Left of original region)
        original_left = 200 * 9525
        original_top = 100 * 9525
        
        width_emu = 600 * 9525
        height_emu = 80 * 9525
        
        # Expansion shifts: Left by 5% width, Top by 1% height
        w_expand = int(width_emu * 0.05)
        h_expand = int(height_emu * 0.01)
        
        expected_left = original_left - w_expand
        expected_top = original_top - h_expand
        
        # Text box should be shifted by expansion
        self.assertEqual(shape.left, expected_left, 
                        "Text box left should be shifted by expansion")
        self.assertEqual(shape.top, expected_top, 
                        "Text box top should be shifted by expansion")


class TestLayoutRect(unittest.TestCase):
    """Unit tests for LayoutRect overlap detection."""

    def test_no_overlap_horizontal(self):
        """Test rects side by side don't overlap."""
        r1 = LayoutRect(left=0, top=0, width=100, height=50, block_index=0)
        r2 = LayoutRect(left=110, top=0, width=100, height=50, block_index=1)
        
        self.assertFalse(r1.overlaps_with(r2))
        self.assertFalse(r2.overlaps_with(r1))
        self.assertEqual(r1.overlap_area(r2), 0)

    def test_no_overlap_vertical(self):
        """Test rects stacked vertically don't overlap."""
        r1 = LayoutRect(left=0, top=0, width=100, height=50, block_index=0)
        r2 = LayoutRect(left=0, top=60, width=100, height=50, block_index=1)
        
        self.assertFalse(r1.overlaps_with(r2))
        self.assertEqual(r1.overlap_area(r2), 0)

    def test_partial_overlap(self):
        """Test partially overlapping rects."""
        r1 = LayoutRect(left=0, top=0, width=100, height=100, block_index=0)
        r2 = LayoutRect(left=50, top=50, width=100, height=100, block_index=1)
        
        self.assertTrue(r1.overlaps_with(r2))
        self.assertTrue(r2.overlaps_with(r1))
        # Overlap: 50x50 = 2500
        self.assertEqual(r1.overlap_area(r2), 2500)

    def test_complete_containment(self):
        """Test when one rect completely contains another."""
        r1 = LayoutRect(left=0, top=0, width=200, height=200, block_index=0)
        r2 = LayoutRect(left=50, top=50, width=50, height=50, block_index=1)
        
        self.assertTrue(r1.overlaps_with(r2))
        # r2 is completely inside r1: 50x50 = 2500
        self.assertEqual(r1.overlap_area(r2), 2500)

    def test_adjacent_no_overlap(self):
        """Test rects that touch edges don't overlap."""
        r1 = LayoutRect(left=0, top=0, width=100, height=100, block_index=0)
        r2 = LayoutRect(left=100, top=0, width=100, height=100, block_index=1)
        
        self.assertFalse(r1.overlaps_with(r2))
        self.assertEqual(r1.overlap_area(r2), 0)

    def test_right_bottom_properties(self):
        """Test right and bottom property calculations."""
        r = LayoutRect(left=10, top=20, width=100, height=50, block_index=0)
        
        self.assertEqual(r.right, 110)
        self.assertEqual(r.bottom, 70)


class TestOverlapResolution(unittest.TestCase):
    """Unit tests for overlap resolution (merge) logic."""

    def setUp(self):
        self.writer = PDFToPPTWriter()
        # Set slide dimensions for testing (large enough to avoid boundary issues)
        self.writer._slide_width = Emu(500000)
        self.writer._slide_height = Emu(500000)

    def test_merge_simple_overlap(self):
        """Test that overlapping rects are merged."""
        # Two rects that overlap
        rects = [
            LayoutRect(left=0, top=0, width=100, height=100, block_index=0, text="A"),
            LayoutRect(left=0, top=50, width=100, height=100, block_index=1, text="B"),
        ]
        
        resolved = self.writer._merge_overlaps(rects)
        
        # Should merge into 1 rect
        self.assertEqual(len(resolved), 1)
        # Check merged dimensions
        self.assertEqual(resolved[0].top, 0) # Min top
        self.assertEqual(resolved[0].height, 150) # Max bottom (150) - Min top (0)
        # Check merged text
        self.assertEqual(resolved[0].text, "A\nB")

    def test_no_change_when_no_overlap(self):
        """Test that non-overlapping rects are not modified."""
        # Place them far apart
        rects = [
            LayoutRect(left=0, top=0, width=100, height=50, block_index=0, text="A"),
            LayoutRect(left=0, top=500000, width=100, height=50, block_index=1, text="B"),
        ]
        
        resolved = self.writer._merge_overlaps(rects)
        
        self.assertEqual(len(resolved), 2)
        self.assertEqual(resolved[0].top, 0)
        self.assertEqual(resolved[1].top, 500000)

    def test_recursive_merge(self):
        """Test that multiple overlapping rects are merged recursively."""
        # Three rects, all at same position
        rects = [
            LayoutRect(left=0, top=0, width=100, height=50, block_index=0, text="A"),
            LayoutRect(left=0, top=0, width=100, height=50, block_index=1, text="B"),
            LayoutRect(left=0, top=0, width=100, height=50, block_index=2, text="C"),
        ]
        
        resolved = self.writer._merge_overlaps(rects)
        
        # Should merge into 1
        self.assertEqual(len(resolved), 1)
        self.assertIn("A", resolved[0].text)
        self.assertIn("B", resolved[0].text)
        self.assertIn("C", resolved[0].text)

    def test_find_overlaps(self):
        """Test _find_overlaps helper method."""
        rects = [
            LayoutRect(left=0, top=0, width=100, height=100, block_index=0),
            LayoutRect(left=50, top=50, width=100, height=100, block_index=1),
            LayoutRect(left=200, top=0, width=50, height=50, block_index=2),
        ]
        
        overlaps = self.writer._find_overlaps(rects)
        
        # Only rect 0 and 1 overlap
        self.assertEqual(len(overlaps), 1)
        self.assertEqual(overlaps[0][0], 0)
        self.assertEqual(overlaps[0][1], 1)
        self.assertEqual(overlaps[0][2], 2500)  # 50x50


class TestDynamicHeight(unittest.TestCase):
    """Unit tests for dynamic height calculation."""

    def setUp(self):
        self.writer = PDFToPPTWriter()
        self.writer._slide_width = Emu(10000000)  # Large enough
        self.writer._slide_height = Emu(10000000)

    def test_single_line_height(self):
        """Test height for single line of text."""
        font_size_emu = int(Pt(12))
        width = int(Pt(200))  # Wide enough for single line
        
        height = self.writer._calculate_dynamic_height("Hello", width, font_size_emu)
        
        # Should be at least font size + padding
        self.assertGreater(height, font_size_emu)

    def test_multiline_height(self):
        """Test that longer text results in greater height."""
        font_size_emu = int(Pt(12))
        narrow_width = int(Pt(50))  # Force wrapping
        
        short_height = self.writer._calculate_dynamic_height("Hi", narrow_width, font_size_emu)
        long_height = self.writer._calculate_dynamic_height(
            "This is a very long text that should wrap to multiple lines", 
            narrow_width, 
            font_size_emu
        )
        
        self.assertGreater(long_height, short_height)

    def test_explicit_newlines(self):
        """Test that explicit newlines increase height."""
        font_size_emu = int(Pt(12))
        width = int(Pt(200))
        
        single_line = self.writer._calculate_dynamic_height("Line 1", width, font_size_emu)
        multi_line = self.writer._calculate_dynamic_height("Line 1\nLine 2\nLine 3", width, font_size_emu)
        
        self.assertGreater(multi_line, single_line)


class TestCJKWidthEstimation(unittest.TestCase):
    """Unit tests for CJK character width estimation."""

    def setUp(self):
        self.writer = PDFToPPTWriter()

    def test_cjk_wider_than_latin(self):
        """Test that CJK text is estimated wider than same-length Latin text."""
        font_size_emu = int(Pt(12))
        
        latin_width = self.writer._estimate_text_width("ABCDEF", font_size_emu)
        cjk_width = self.writer._estimate_text_width("가나다라마바", font_size_emu)
        
        # Same character count, but CJK should be wider
        self.assertGreater(cjk_width, latin_width)

    def test_is_cjk_char(self):
        """Test CJK character detection."""
        # Korean
        self.assertTrue(self.writer._is_cjk_char('한'))
        self.assertTrue(self.writer._is_cjk_char('글'))
        
        # Chinese
        self.assertTrue(self.writer._is_cjk_char('中'))
        self.assertTrue(self.writer._is_cjk_char('文'))
        
        # Japanese (Kanji)
        self.assertTrue(self.writer._is_cjk_char('日'))
        
        # Latin
        self.assertFalse(self.writer._is_cjk_char('A'))
        self.assertFalse(self.writer._is_cjk_char('z'))
        
        # Numbers
        self.assertFalse(self.writer._is_cjk_char('1'))

    def test_mixed_text_width(self):
        """Test width estimation for mixed CJK and Latin text."""
        font_size_emu = int(Pt(12))
        
        # Mixed text should be between pure Latin and pure CJK
        latin_width = self.writer._estimate_text_width("ABCD", font_size_emu)
        cjk_width = self.writer._estimate_text_width("가나다라", font_size_emu)
        mixed_width = self.writer._estimate_text_width("AB가나", font_size_emu)
        
        self.assertGreater(mixed_width, latin_width)
        self.assertLess(mixed_width, cjk_width)


class TestAdaptiveColor(unittest.TestCase):
    """Unit tests for adaptive color logic."""

    def setUp(self):
        self.writer = PDFToPPTWriter()

    def test_get_luminance(self):
        """Test luminance calculation."""
        # White = 1.0
        self.assertAlmostEqual(self.writer._get_luminance((255, 255, 255)), 1.0)
        # Black = 0.0
        self.assertAlmostEqual(self.writer._get_luminance((0, 0, 0)), 0.0)
        # Red (0.299)
        self.assertAlmostEqual(self.writer._get_luminance((255, 0, 0)), 0.299, places=2)
        # Green (0.587)
        self.assertAlmostEqual(self.writer._get_luminance((0, 255, 0)), 0.587, places=2)

    def test_get_optimal_text_color(self):
        """Test text color selection based on background."""
        # Bright BG -> Black Text
        self.assertEqual(self.writer._get_optimal_text_color((255, 255, 255)), (0, 0, 0))
        # Dark BG -> White Text
        self.assertEqual(self.writer._get_optimal_text_color((0, 0, 0)), (255, 255, 255))
        # Mid-gray -> Check threshold
        # Gray 128 is ~0.5. 
        # (128/255)*1.0 ~= 0.501 -> Bright -> Black
        self.assertEqual(self.writer._get_optimal_text_color((128, 128, 128)), (0, 0, 0))

    def test_get_average_color(self):
        """Test average color extraction from image region."""
        if not Image:
            self.skipTest("PIL not installed")
            
        width, height = 100, 100
        # Create a red image
        image = Image.new("RGB", (width, height), "red")
        
        rect = LayoutRect(left=0, top=0, width=500000, height=500000, block_index=0)
        # Scale 10000 means 500000 emu = 50 px
        scale = 10000.0
        
        avg_color = self.writer._get_average_color(image, rect, scale, 0, 0)
        self.assertEqual(avg_color, (255, 0, 0))


if __name__ == "__main__":
    unittest.main()
