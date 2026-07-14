"""Unit tests for the review/edit session (WP-C5)."""

from __future__ import annotations

import io
import json
import tempfile
import unittest
from copy import deepcopy
from pathlib import Path
from unittest import mock

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from src.chains.color_distribution_chain import ColoredSegment
from src.core.ppt_parser import PPTParser
from src.core.ppt_writer import snapshot_fit_geometry
from src.services.consistency_sweep import Finding
from src.services.quality_records import QualityRecorder
from src.services.review_session import ReviewSession


def _deck(text: str, n: int) -> io.BytesIO:
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for _ in range(n):
        s = prs.slides.add_slide(blank)
        tb = s.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
        tb.text_frame.text = text
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def _session(text: str, n: int, targets=None, findings=None) -> ReviewSession:
    paras, pres = PPTParser().extract_paragraphs(_deck(text, n))
    targets = targets or [f"T{i}" for i in range(len(paras))]
    return ReviewSession(
        presentation=pres, paragraphs=paras, translated_texts=targets,
        findings=findings or [], source_lang="한국어", target_lang="영어",
        model="stub",
    )


def _multicolor_session() -> ReviewSession:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    paragraph = box.text_frame.paragraphs[0]
    red = paragraph.add_run()
    red.text = "강력히 "
    red.font.color.rgb = RGBColor(255, 0, 0)
    blue = paragraph.add_run()
    blue.text = "추천합니다"
    blue.font.color.rgb = RGBColor(0, 0, 255)
    source = io.BytesIO()
    prs.save(source)
    source.seek(0)
    paragraphs, presentation = PPTParser().extract_paragraphs(source)
    target = "This approach is strongly recommended"
    distribution = {
        0: [
            ColoredSegment(text="This approach is ", group_index=1),
            ColoredSegment(text="strongly", group_index=0),
            ColoredSegment(text=" recommended", group_index=1),
        ]
    }
    return ReviewSession(
        presentation=presentation,
        paragraphs=paragraphs,
        translated_texts=[target],
        findings=[],
        source_lang="한국어",
        target_lang="영어",
        model="stub",
        color_distributions=distribution,
        source_pptx=source.getvalue(),
    )


class FragmentsTestCase(unittest.TestCase):
    def test_repeat_count_reflects_duplicates(self) -> None:
        sess = _session("자세한 내용은 부록 참조", 4)
        frags = sess.fragments()
        self.assertEqual(len(frags), 4)
        self.assertTrue(all(f.repeat_count == 4 for f in frags))

    def test_findings_attach_to_fragments(self) -> None:
        finding = Finding(
            type="terminology.violation", severity="major", description="x",
            location={"slide": 1, "shape": 0, "paragraph": 0},
            segment={"source": "s", "output": "o"}, ordinal=1,
            fragment_index=0, suggested_fix="Aim Punch",
        )
        sess = _session("A", 2, findings=[finding])
        frags = sess.fragments()
        self.assertEqual(len(frags[0].findings), 1)
        self.assertEqual(frags[0].findings[0]["type"], "terminology.violation")
        self.assertEqual(frags[1].findings, [])


class EditPropagationTestCase(unittest.TestCase):
    def test_identical_indices(self) -> None:
        sess = _session("반복", 3)
        self.assertEqual(sorted(sess.identical_indices(0)), [0, 1, 2])

    def test_edit_propagates_to_identicals(self) -> None:
        sess = _session("반복", 3)
        changed = sess.apply_edit(0, "EDITED", propagate_identical=True)
        self.assertEqual(sorted(changed), [0, 1, 2])
        self.assertTrue(all(t == "EDITED" for t in sess.translated_texts))

    def test_edit_without_propagation_only_target(self) -> None:
        sess = _session("반복", 3)
        changed = sess.apply_edit(0, "EDITED", propagate_identical=False)
        self.assertEqual(changed, [0])
        self.assertEqual(sess.translated_texts[1], "T1")

    def test_render_reflects_edit(self) -> None:
        sess = _session("반복", 2, targets=["OLD", "OLD"])
        sess.apply_edit(0, "NEW", propagate_identical=True)
        out = sess.render()
        check = Presentation(out)
        texts = [sh.text_frame.text for sl in check.slides for sh in sl.shapes if sh.has_text_frame]
        self.assertTrue(all(t == "NEW" for t in texts))


class IdempotentRenderTestCase(unittest.TestCase):
    """C-2: repeated render() must not cumulatively shrink fonts."""

    def _fitting_deck(self) -> io.BytesIO:
        prs = Presentation()
        s = prs.slides.add_slide(prs.slide_layouts[6])
        tb = s.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        run = tb.text_frame.paragraphs[0].add_run()
        run.text = "짧은"
        run.font.size = Pt(20)
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf

    def _font_size_pt(self, buffer: io.BytesIO):
        prs = Presentation(buffer)
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size is not None:
                            return run.font.size.pt
        return None

    def test_repeated_render_is_idempotent_with_snapshot(self) -> None:
        paras, pres = PPTParser().extract_paragraphs(self._fitting_deck())
        # A long translation forces auto_shrink to reduce the font.
        long_target = "This is a considerably longer translated sentence" * 2
        snapshot = snapshot_fit_geometry(pres)
        sess = ReviewSession(
            presentation=pres, paragraphs=paras, translated_texts=[long_target],
            findings=[], source_lang="한국어", target_lang="영어", model="stub",
            text_fit_mode="auto_shrink", min_font_ratio=50, fit_snapshot=snapshot,
        )

        size1 = self._font_size_pt(sess.render())
        size2 = self._font_size_pt(sess.render())
        size3 = self._font_size_pt(sess.render())

        self.assertIsNotNone(size1)
        # Fit shrank below the original 20pt...
        self.assertLess(size1, 20)
        # ...but repeated renders reproduce the same size, not 0.x^k of it.
        self.assertAlmostEqual(size1, size2, places=3)
        self.assertAlmostEqual(size2, size3, places=3)

    def _big_font_deck(self) -> io.BytesIO:
        prs = Presentation()
        s = prs.slides.add_slide(prs.slide_layouts[6])
        tb = s.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        run = tb.text_frame.paragraphs[0].add_run()
        run.text = "ab"  # 2 chars
        run.font.size = Pt(40)
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf

    def test_render_without_snapshot_is_still_idempotent(self) -> None:
        # Fresh-source rendering no longer depends on an optional geometry
        # snapshot; every render starts from the pristine source bytes.
        paras, pres = PPTParser().extract_paragraphs(self._big_font_deck())
        sess = ReviewSession(
            presentation=pres, paragraphs=paras, translated_texts=["abcd"],
            findings=[], source_lang="한국어", target_lang="영어", model="stub",
            text_fit_mode="auto_shrink", min_font_ratio=50, fit_snapshot=None,
        )
        size1 = self._font_size_pt(sess.render())
        size2 = self._font_size_pt(sess.render())
        self.assertAlmostEqual(size2, size1, places=3)


class TransactionalColorRenderTestCase(unittest.TestCase):
    @staticmethod
    def _runs(buffer: io.BytesIO):
        presentation = Presentation(buffer)
        paragraph = presentation.slides[0].shapes[0].text_frame.paragraphs[0]
        return [
            (run.text, str(run.font.color.rgb) if run.font.color.rgb else None)
            for run in paragraph.runs
            if run.text
        ]

    def test_reordered_color_mapping_is_stable_across_renders(self) -> None:
        session = _multicolor_session()
        first = self._runs(session.render())
        second = self._runs(session.render())
        self.assertEqual(first, second)
        self.assertEqual(
            first,
            [
                ("This approach is ", "0000FF"),
                ("strongly", "FF0000"),
                (" recommended", "0000FF"),
            ],
        )

    def test_same_text_save_is_noop_and_keeps_color_distribution(self) -> None:
        session = _multicolor_session()
        before = deepcopy(session.color_distributions)
        changed = session.apply_edit(0, session.translated_texts[0])
        self.assertEqual(changed, [])
        self.assertEqual(session.color_distributions, before)
        self.assertEqual(session.revision, 0)


class DraftProposalTestCase(unittest.TestCase):
    def test_proposal_does_not_mutate_until_apply_and_can_undo(self) -> None:
        session = _session("원문", 1, targets=["OLD"])
        proposal = session.create_proposal(
            0,
            action="edit",
            target="NEW",
            instruction=None,
            propagate_identical=False,
            model="stub",
            provider="anthropic",
        )
        self.assertEqual(session.translated_texts[0], "OLD")
        self.assertEqual(session.revision, 0)

        applied = session.apply_proposal(proposal.id, expected_revision=0)
        self.assertEqual(applied.changed_indices, [0])
        self.assertEqual(session.translated_texts[0], "NEW")
        self.assertTrue(session.dirty)
        self.assertEqual(session.revision, 1)

        undone = session.undo(expected_revision=1)
        self.assertEqual(undone, [0])
        self.assertEqual(session.translated_texts[0], "OLD")
        self.assertEqual(session.revision, 2)

    def test_stale_proposal_is_rejected(self) -> None:
        session = _session("원문", 1, targets=["OLD"])
        proposal = session.create_proposal(
            0,
            action="edit",
            target="NEW",
            instruction=None,
            propagate_identical=False,
            model="stub",
            provider="anthropic",
        )
        session.apply_edit(0, "OTHER")
        with self.assertRaises(RuntimeError):
            session.apply_proposal(proposal.id, expected_revision=0)

    def test_partial_candidates_use_changed_phrase(self) -> None:
        session = _session("원문", 3, targets=[
            "Adjusted field drop rates",
            "The field drop table",
            "Unrelated text",
        ])
        proposal = session.create_proposal(
            0,
            action="edit",
            target="Adjusted World Spawn rates",
            instruction=None,
            propagate_identical=False,
            model="stub",
            provider="anthropic",
        )
        self.assertEqual(len(proposal.partial_candidates), 1)
        candidate = proposal.partial_candidates[0]
        self.assertEqual(candidate["old_phrase"], "field drop")
        self.assertEqual(candidate["new_phrase"], "World Spawn")
        self.assertEqual(candidate["proposed_target"], "The World Spawn table")


class PartialMatchTestCase(unittest.TestCase):
    def test_partial_candidates_found_but_not_auto_applied(self) -> None:
        sess = _session("A", 1, targets=["World Spawn item list"])
        # Manually add a second fragment target containing the phrase.
        sess.translated_texts.append("Adjusted field drop rates")
        # Simulate an edit that introduces "World Spawn".
        candidates = sess.partial_match_candidates(0, "World Spawn")
        # No fragment (other than index 0) contains "World Spawn" yet.
        self.assertEqual(candidates, [])

    def test_single_hangul_character_change_is_not_suggested_across_deck(self) -> None:
        sess = _session("원문", 4, targets=[
            "활성화 기준",
            "3개월 개발 + 6주 안정화, 11월 중순 베타 테스트",
            "Phase 2 범위는 활성화, 명확성, 매치메이킹에 초점",
            "원인이 방향성 실패인지 활성화 실패인지가 핵심 질문",
        ])
        proposal = sess.create_proposal(
            0,
            action="edit",
            target="활성총 기준",
            instruction=None,
            propagate_identical=False,
            model="stub",
            provider="anthropic",
        )
        self.assertEqual(ReviewSession._changed_phrase("활성화 기준", "활성총 기준"), ("화", "총"))
        self.assertEqual(proposal.partial_candidates, [])

    def test_common_two_character_cjk_phrase_is_not_suggested_across_deck(self) -> None:
        sess = _session("원문", 3, targets=[
            "활성화 기준",
            "품질 기준",
            "출시 기준",
        ])
        proposal = sess.create_proposal(
            0,
            action="edit",
            target="활성화 규칙",
            instruction=None,
            propagate_identical=False,
            model="stub",
            provider="anthropic",
        )
        self.assertEqual(ReviewSession._changed_phrase("활성화 기준", "활성화 규칙"), ("기준", "규칙"))
        self.assertEqual(proposal.partial_candidates, [])

    def test_short_lowercase_latin_word_is_not_suggested_across_deck(self) -> None:
        sess = _session("source", 3, targets=[
            "replace the policy",
            "the schedule",
            "review the results",
        ])
        proposal = sess.create_proposal(
            0,
            action="edit",
            target="replace a policy",
            instruction=None,
            propagate_identical=False,
            model="stub",
            provider="anthropic",
        )
        self.assertEqual(ReviewSession._changed_phrase("replace the policy", "replace a policy"), ("the", "a"))
        self.assertEqual(proposal.partial_candidates, [])

    def test_latin_partial_match_requires_word_boundaries(self) -> None:
        sess = _session("source", 3, targets=[
            "CAT policy",
            "CONCAT function",
            "The CAT schedule",
        ])
        proposal = sess.create_proposal(
            0,
            action="edit",
            target="DOG policy",
            instruction=None,
            propagate_identical=False,
            model="stub",
            provider="anthropic",
        )
        self.assertEqual(
            [candidate["target"] for candidate in proposal.partial_candidates],
            ["The CAT schedule"],
        )

    def test_accented_latin_partial_match_requires_unicode_word_boundaries(self) -> None:
        sess = _session("source", 3, targets=[
            "café policy",
            "caféteria",
            "Le café schedule",
        ])
        proposal = sess.create_proposal(
            0,
            action="edit",
            target="bistro policy",
            instruction=None,
            propagate_identical=False,
            model="stub",
            provider="anthropic",
        )
        self.assertEqual(
            [candidate["target"] for candidate in proposal.partial_candidates],
            ["Le café schedule"],
        )

    def test_hangul_phrase_inside_longer_word_is_not_suggested(self) -> None:
        sess = _session("원문", 3, targets=[
            "Phase 2 활성화 기준",
            "비활성화 정책",
            "활성화 일정",
        ])
        proposal = sess.create_proposal(
            0,
            action="edit",
            target="Phase 2 기준",
            instruction=None,
            propagate_identical=False,
            model="stub",
            provider="anthropic",
        )
        self.assertEqual(
            [(candidate["target"], candidate["proposed_target"])
             for candidate in proposal.partial_candidates],
            [("활성화 일정", "일정")],
        )

    def test_japanese_phrase_without_spaces_remains_supported(self) -> None:
        sess = _session("source", 2, targets=[
            "有効化 policy",
            "新規有効化基準",
        ])
        proposal = sess.create_proposal(
            0,
            action="edit",
            target="activation policy",
            instruction=None,
            propagate_identical=False,
            model="stub",
            provider="anthropic",
        )
        self.assertEqual(
            [(candidate["target"], candidate["proposed_target"])
             for candidate in proposal.partial_candidates],
            [("新規有効化基準", "新規activation基準")],
        )

    def test_deletion_preview_matches_applied_result(self) -> None:
        sess = _session("원문", 2, targets=["Phase 2 활성화 기준", "활성화 일정"])
        proposal = sess.create_proposal(
            0,
            action="edit",
            target="Phase 2 기준",
            instruction=None,
            propagate_identical=False,
            model="stub",
            provider="anthropic",
        )
        self.assertEqual(len(proposal.partial_candidates), 1)
        candidate = proposal.partial_candidates[0]
        self.assertEqual(candidate["old_phrase"], "활성화")
        self.assertEqual(candidate["new_phrase"], "")
        self.assertEqual(candidate["proposed_target"], "일정")

        sess.apply_proposal(proposal.id, expected_revision=0)
        changed = sess.apply_partial_candidates(
            [candidate["index"]],
            old_phrase=candidate["old_phrase"],
            new_phrase=candidate["new_phrase"],
            expected_revision=1,
            model="stub",
            provider="anthropic",
        )
        self.assertEqual(changed, [1])
        self.assertEqual(sess.translated_texts[1], candidate["proposed_target"])

    def test_identical_source_propagation_is_not_repeated_as_partial_candidate(self) -> None:
        sess = _session("동일 원문", 3, targets=["Phase 2 활성화 기준"] * 3)
        proposal = sess.create_proposal(
            0,
            action="edit",
            target="Phase 2 운영 기준",
            instruction=None,
            propagate_identical=True,
            model="stub",
            provider="anthropic",
        )
        self.assertEqual(proposal.changed_indices, [0, 1, 2])
        self.assertEqual(proposal.partial_candidates, [])

    def test_ambiguous_repeated_phrase_in_one_target_is_not_suggested(self) -> None:
        sess = _session("source", 2, targets=[
            "field drop policy",
            "field drop follows field drop rules",
        ])
        proposal = sess.create_proposal(
            0,
            action="edit",
            target="World Spawn policy",
            instruction=None,
            propagate_identical=False,
            model="stub",
            provider="anthropic",
        )
        self.assertEqual(proposal.partial_candidates, [])


class PartialMatchFixtureRegressionTestCase(unittest.TestCase):
    """End-to-end candidate checks using the synthetic false-positive deck."""

    @classmethod
    def setUpClass(cls) -> None:
        fixture = Path(__file__).parent / "fixtures" / "partial-match-false-positive.pptx"
        cls.source_pptx = fixture.read_bytes()
        cls.paragraphs, cls.presentation = PPTParser().extract_paragraphs(
            io.BytesIO(cls.source_pptx)
        )
        cls.scenario_texts = {
            "활성화 기준",
            "3개월 개발 + 6주 안정화, 11월 중순 베타 테스트",
            "Phase 2 범위는 활성화, 명확성, 매치메이킹에 초점",
            "원인이 방향성 실패인지 활성화 실패인지가 핵심 질문",
            "전화 회의 결과를 문서화하고 시각화 자료를 정리",
            "명확화 작업과 안정화 작업은 서로 다른 일정으로 관리",
            "CAT policy",
            "CONCAT function",
            "The CAT schedule",
            "SCATTER plot",
            "CAT-based workflow",
            "Category policy",
            "Phase 2 활성화 기준",
            "활성화 일정",
            "다음 단계 활성화",
            "비활성화 정책",
            "활성화",
            "활성화 여부와 활성화 일정",
            "Adjusted field drop rates",
            "The field drop table",
            "field drop follows field drop rules",
            "Unrelated text",
        }

    def _session(self) -> ReviewSession:
        targets = [
            text if (text := (paragraph.original_text or "").strip()) in self.scenario_texts else ""
            for paragraph in self.paragraphs
        ]
        return ReviewSession(
            presentation=deepcopy(self.presentation),
            paragraphs=self.paragraphs,
            translated_texts=targets,
            findings=[],
            source_lang="한국어",
            target_lang="한국어",
            model="stub",
            source_pptx=self.source_pptx,
        )

    def _index(self, text: str, occurrence: int = 0) -> int:
        matches = [
            index
            for index, paragraph in enumerate(self.paragraphs)
            if (paragraph.original_text or "").strip() == text
        ]
        return matches[occurrence]

    def test_fixture_blocks_false_positives_and_keeps_valid_candidates(self) -> None:
        session = self._session()
        hangul = session.create_proposal(
            self._index("활성화 기준"),
            action="edit",
            target="활성총 기준",
            instruction=None,
            propagate_identical=False,
            model="stub",
            provider="anthropic",
        )
        self.assertEqual(hangul.partial_candidates, [])

        latin = session.create_proposal(
            self._index("CAT policy"),
            action="edit",
            target="DOG policy",
            instruction=None,
            propagate_identical=False,
            model="stub",
            provider="anthropic",
        )
        self.assertEqual(
            [candidate["target"] for candidate in latin.partial_candidates],
            ["The CAT schedule", "CAT-based workflow"],
        )

        valid_phrase = session.create_proposal(
            self._index("Adjusted field drop rates"),
            action="edit",
            target="Adjusted World Spawn rates",
            instruction=None,
            propagate_identical=False,
            model="stub",
            provider="anthropic",
        )
        self.assertEqual(
            [candidate["target"] for candidate in valid_phrase.partial_candidates],
            ["The field drop table"],
        )

        deletion = session.create_proposal(
            self._index("Phase 2 활성화 기준"),
            action="edit",
            target="Phase 2 기준",
            instruction=None,
            propagate_identical=True,
            model="stub",
            provider="anthropic",
        )
        deletion_by_target = {
            candidate["target"]: candidate["proposed_target"]
            for candidate in deletion.partial_candidates
        }
        self.assertNotIn("비활성화 정책", deletion_by_target)
        self.assertNotIn("활성화 여부와 활성화 일정", deletion_by_target)
        self.assertEqual(deletion_by_target["활성화 일정"], "일정")
        self.assertEqual(deletion_by_target["다음 단계 활성화"], "다음 단계")

        identical = session.create_proposal(
            self._index("Phase 2 활성화 기준"),
            action="edit",
            target="Phase 2 운영 기준",
            instruction=None,
            propagate_identical=True,
            model="stub",
            provider="anthropic",
        )
        self.assertEqual(len(identical.changed_indices), 3)
        self.assertTrue(
            set(identical.changed_indices).isdisjoint(
                candidate["index"] for candidate in identical.partial_candidates
            )
        )


class LengthBudgetTestCase(unittest.TestCase):
    def test_note_has_no_budget(self) -> None:
        sess = _session("본문", 1)
        # No notes in this deck; body fragment has a budget.
        self.assertIsNotNone(sess.length_budget(0))

    def test_length_limit_violation_is_visible_in_review(self) -> None:
        sess = _session("1234567890", 1, targets=["x" * 14])
        sess.length_limit = 130

        findings = sess.run_final_sweep()

        length_findings = [f for f in findings if f.type == "fit.length_limit"]
        self.assertEqual(len(length_findings), 1)
        self.assertEqual(length_findings[0].fragment_index, 0)
        self.assertIn("최대 13자", length_findings[0].description)


class RetranslateFragmentTestCase(unittest.TestCase):
    """A-1: retranslate logic lives on ReviewSession, chain mocked out."""

    def test_returns_new_target_and_threads_instruction_and_budget(self) -> None:
        sess = _session("반복", 1, targets=["OLD"])
        sess.length_limit = 130

        captured: dict = {}

        def fake_create_chain(**kwargs):
            captured["instructions"] = kwargs.get("instructions")
            captured["model_name"] = kwargs.get("model_name")
            captured["provider"] = kwargs.get("provider")
            captured["length_limit"] = kwargs.get("length_limit")
            return object()  # opaque chain handle; translate is mocked below

        with mock.patch(
            "src.chains.translation_chain.create_translation_chain",
            side_effect=fake_create_chain,
        ), mock.patch(
            "src.chains.translation_chain.translate_with_progress",
            return_value=["NEW"],
        ):
            new_target, color_segments = sess.retranslate_fragment(
                0, "격식체로", model="m", provider="anthropic"
            )

        self.assertEqual(new_target, "NEW")
        # Single-color fragment -> no color segments re-mapped.
        self.assertIsNone(color_segments)
        self.assertEqual(captured["model_name"], "m")
        self.assertEqual(captured["provider"], "anthropic")
        self.assertEqual(captured["length_limit"], 130)
        # Both the user instruction and the length-budget guidance are threaded in.
        self.assertIn("격식체로", captured["instructions"])
        self.assertIn("슬라이드 박스", captured["instructions"])

    def test_does_not_mutate_session_state(self) -> None:
        sess = _session("반복", 1, targets=["OLD"])
        with mock.patch(
            "src.chains.translation_chain.create_translation_chain",
            return_value=object(),
        ), mock.patch(
            "src.chains.translation_chain.translate_with_progress",
            return_value=["NEW"],
        ):
            sess.retranslate_fragment(0, None, model="m", provider="anthropic")
        # The caller (api.py) applies the edit; the method itself is pure.
        self.assertEqual(sess.translated_texts[0], "OLD")

    def test_empty_result_raises(self) -> None:
        sess = _session("반복", 1, targets=["OLD"])
        with mock.patch(
            "src.chains.translation_chain.create_translation_chain",
            return_value=object(),
        ), mock.patch(
            "src.chains.translation_chain.translate_with_progress",
            return_value=[],
        ):
            with self.assertRaises(RuntimeError):
                sess.retranslate_fragment(0, None, model="m", provider="anthropic")


class RecordReviewEditTestCase(unittest.TestCase):
    def test_accepted_edit_records_triplet(self) -> None:
        with tempfile.TemporaryDirectory() as d:
            rec = QualityRecorder(quality_dir=d)
            rec.record_review_edit(
                job_id="deck:x.pptx", doc_ref="deck:x.pptx",
                source_lang="한국어", target_lang="영어", model="stub",
                location={"slide": 3, "shape": 1, "paragraph": 0},
                segment={"source": "원문", "output": "old", "corrected": "new"},
                disposition="accepted", propagated=3,
            )
            rows = [json.loads(l) for l in (Path(d) / "quality_records.jsonl").read_text().splitlines()]
            self.assertEqual(len(rows), 1)
            r = rows[0]
            self.assertEqual(r["origin"]["stage"], "manual_edit")
            self.assertEqual(r["origin"]["executor"], "human")
            self.assertEqual(r["disposition"], "accepted")
            self.assertEqual(r["segment"]["corrected"], "new")
            self.assertEqual(r["promotion"]["status"], "candidate")

    def test_rejected_ignore_records(self) -> None:
        with tempfile.TemporaryDirectory() as d:
            rec = QualityRecorder(quality_dir=d)
            rec.record_review_edit(
                job_id="j", doc_ref="deck:x.pptx",
                source_lang="한국어", target_lang="영어", model="stub",
                location={"slide": 1, "shape": 0, "paragraph": 0},
                segment={"source": "원문", "output": "out", "corrected": None},
                disposition="rejected", finding_type="consistency.phrase",
            )
            rows = [json.loads(l) for l in (Path(d) / "quality_records.jsonl").read_text().splitlines()]
            self.assertEqual(rows[0]["disposition"], "rejected")
            self.assertEqual(rows[0]["finding"]["type"], "consistency.phrase")
            self.assertEqual(rows[0]["promotion"]["status"], "rejected")


if __name__ == "__main__":
    unittest.main()
