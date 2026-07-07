"""Tests for src/core/text_extractor.py (P1 improvements).

Covers reading order, title dedup, table cell newline escaping, and chart
data extraction, based on the reproduction deck from
docs/MD_EXTRACTION_IMPROVEMENT_PLAN.md.
"""

from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from src.core.text_extractor import (
    ChartBlock,
    ExtractionOptions,
    TextBlock,
    docs_to_markdown,
    extract_pptx_to_docs,
)


@pytest.fixture(scope="module")
def demo_pptx(tmp_path_factory) -> str:
    """Build the reproduction deck from the improvement plan."""
    prs = Presentation()

    # Slide 1: title placeholder + body + multi-paragraph note
    s1 = prs.slides.add_slide(prs.slide_layouts[1])
    s1.shapes.title.text = "2026 사업 전략"
    body = s1.placeholders[1].text_frame
    body.text = "시장 규모 12% 성장"
    p = body.add_paragraph()
    p.text = "경쟁 심화"
    p.level = 1
    s1.notes_slide.notes_text_frame.text = "첫 문단입니다.\n둘째 문단: 발표 시 강조."

    # Slide 2: z-order trap (right column inserted first, heading last)
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    s2.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4), Inches(3)).text_frame.text = (
        "[오른쪽 칼럼] 하반기 계획"
    )
    s2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(3)).text_frame.text = (
        "[왼쪽 칼럼] 상반기 실적"
    )
    s2.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(8), Inches(0.8)).text_frame.text = (
        "연간 로드맵"
    )

    # Slide 3: table with in-cell newline + merged cell, plus a chart
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    tbl = s3.shapes.add_table(3, 3, Inches(0.5), Inches(0.5), Inches(5), Inches(2)).table
    for c, h in enumerate(["구분", "2025", "2026"]):
        tbl.cell(0, c).text = h
    tbl.cell(1, 0).text = "매출\n(단위: 억)"
    tbl.cell(1, 1).text = "120"
    tbl.cell(1, 2).text = "150"
    tbl.cell(2, 0).merge(tbl.cell(2, 2))
    tbl.cell(2, 0).text = "출처: 내부 집계"
    cd = CategoryChartData()
    cd.categories = ["Q1", "Q2", "Q3"]
    cd.add_series("매출", (30, 40, 50))
    cd.add_series("이익", (5, 8, 12))
    s3.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(3), Inches(5), Inches(3), cd
    )

    path = str(tmp_path_factory.mktemp("extractor") / "demo.pptx")
    prs.save(path)
    return path


@pytest.fixture(scope="module")
def options() -> ExtractionOptions:
    return ExtractionOptions(
        figures="placeholder", charts="labels", table_header=True, with_notes=True
    )


@pytest.fixture(scope="module")
def docs(demo_pptx, options):
    return extract_pptx_to_docs(demo_pptx, options)


@pytest.fixture(scope="module")
def markdown(docs, options) -> str:
    return docs_to_markdown(docs, options)


# --- P1-1: reading order ---------------------------------------------------


def test_slide2_title_is_topmost_text(docs):
    assert docs[1].title == "연간 로드맵"


def test_slide2_reading_order_left_before_right(docs):
    text_blocks = [b for b in docs[1].blocks if isinstance(b, TextBlock)]
    flattened = [line for b in text_blocks for line in b.lines]
    assert flattened.index("[왼쪽 칼럼] 상반기 실적") < flattened.index("[오른쪽 칼럼] 하반기 계획")


def test_slide2_markdown_order(markdown):
    slide2 = markdown.split("## Slide 2")[1].split("## Slide 3")[0]
    assert slide2.index("연간 로드맵") < slide2.index("[왼쪽 칼럼]") < slide2.index("[오른쪽 칼럼]")


# --- P1-2: title dedup -----------------------------------------------------


def test_slide1_title_appears_once(markdown):
    assert markdown.count("2026 사업 전략") == 1


def test_slide2_fallback_title_not_repeated_as_bullet(markdown):
    slide2 = markdown.split("## Slide 2")[1].split("## Slide 3")[0]
    # Appears once, in the heading only; never in the body below it.
    assert slide2.count("연간 로드맵") == 1
    body = "\n".join(slide2.splitlines()[1:])  # drop heading remainder
    assert "연간 로드맵" not in body


def test_fallback_title_shape_keeps_remaining_lines(tmp_path, options):
    """Only the first line of a fallback title shape is consumed."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tf = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4), Inches(2)).text_frame
    tf.text = "제목 줄"
    p = tf.add_paragraph()
    p.text = "본문 줄"
    path = str(tmp_path / "fallback.pptx")
    prs.save(path)

    doc = extract_pptx_to_docs(path, options)[0]
    assert doc.title == "제목 줄"
    lines = [line for b in doc.blocks if isinstance(b, TextBlock) for line in b.lines]
    assert "제목 줄" not in lines
    assert "본문 줄" in lines


# --- P1-3: table cell newline escaping --------------------------------------


def test_table_cell_newline_becomes_br(markdown):
    assert "매출<br>(단위: 억)" in markdown


def test_table_rows_are_valid_markdown(markdown):
    slide3 = markdown.split("## Slide 3")[1]
    table_lines = [l for l in slide3.splitlines() if l.startswith("|")]
    assert table_lines, "table missing from slide 3"
    for line in table_lines:
        assert line.startswith("|") and line.endswith("|")
        assert "\n" not in line


# --- P1-4: chart data extraction --------------------------------------------


def test_chart_block_extracted(docs):
    charts = [b for b in docs[2].blocks if isinstance(b, ChartBlock)]
    assert len(charts) == 1
    chart = charts[0]
    assert chart.categories == ["Q1", "Q2", "Q3"]
    names = [name for name, _ in chart.series]
    assert names == ["매출", "이익"]
    assert [list(v) for _, v in chart.series] == [[30.0, 40.0, 50.0], [5.0, 8.0, 12.0]]


def test_chart_rendered_as_table(markdown):
    slide3 = markdown.split("## Slide 3")[1]
    assert "| 구분 | Q1 | Q2 | Q3 |" in slide3
    assert "| 매출 | 30 | 40 | 50 |" in slide3
    assert "| 이익 | 5 | 8 | 12 |" in slide3


# --- regression guards -------------------------------------------------------


def test_slide1_body_and_note_preserved(markdown):
    slide1 = markdown.split("## Slide 2")[0]
    assert "- 시장 규모 12% 성장" in slide1
    assert "  - 경쟁 심화" in slide1
    assert "> NOTE:" in slide1


def test_slide_count(docs):
    assert len(docs) == 3
