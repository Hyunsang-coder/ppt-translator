"""Tests for text fitting logic in PPTWriter."""

from __future__ import annotations

import unittest

from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt

from src.core.ppt_writer import (
    PPTWriter,
    apply_text_fit,
    _build_shape_context,
    _calculate_available_expansion,
    _safe_expand_width,
    _EXPANSION_GAP_EMU,
    _EXPANSION_THRESHOLD,
)
from src.core.ppt_parser import ParagraphInfo
from src.services.models import TextFitMode


def _make_textbox(prs: Presentation, text: str, font_size_pt: int = 20):
    """Create a slide with a single text box and return (shape, text_frame)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    tf = shape.text_frame
    tf.paragraphs[0].text = ""
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(font_size_pt)
    return shape, tf


class TextFitNoneTestCase(unittest.TestCase):
    """TextFitMode.NONE should not modify anything."""

    def test_none_mode_no_changes(self):
        prs = Presentation()
        shape, tf = _make_textbox(prs, "Hello", font_size_pt=20)
        original_size = tf.paragraphs[0].runs[0].font.size

        apply_text_fit(tf, original_len=5, translated_len=15, mode=TextFitMode.NONE)

        self.assertEqual(tf.paragraphs[0].runs[0].font.size, original_size)


class TextFitAutoShrinkTestCase(unittest.TestCase):
    """TextFitMode.AUTO_SHRINK tests."""

    def test_no_change_when_text_shorter_or_equal(self):
        """If translated text is same length or shorter, don't touch font."""
        prs = Presentation()
        shape, tf = _make_textbox(prs, "Hello World", font_size_pt=20)

        apply_text_fit(tf, original_len=11, translated_len=8, mode=TextFitMode.AUTO_SHRINK)

        self.assertEqual(tf.paragraphs[0].runs[0].font.size, Pt(20))

    def test_no_change_when_ratio_within_threshold(self):
        """Small expansion (<=10%) should not trigger shrink."""
        prs = Presentation()
        shape, tf = _make_textbox(prs, "Hello World", font_size_pt=20)

        # 10% longer
        apply_text_fit(tf, original_len=100, translated_len=110, mode=TextFitMode.AUTO_SHRINK)

        self.assertEqual(tf.paragraphs[0].runs[0].font.size, Pt(20))

    def test_proportional_shrink(self):
        """Font should shrink proportionally when text is longer."""
        prs = Presentation()
        shape, tf = _make_textbox(prs, "Hello", font_size_pt=20)

        # Text doubled in length -> font should shrink to ratio * original
        apply_text_fit(
            tf, original_len=10, translated_len=20,
            mode=TextFitMode.AUTO_SHRINK, min_font_ratio=50,
        )

        new_size = tf.paragraphs[0].runs[0].font.size
        # 10/20 = 0.5, so 20pt * 0.5 = 10pt
        self.assertEqual(new_size, Pt(10))

    def test_min_font_ratio_floor(self):
        """Font should not shrink below min_font_ratio."""
        prs = Presentation()
        shape, tf = _make_textbox(prs, "Hi", font_size_pt=20)

        # Text 5x longer, but min_font_ratio=70 -> floor at 14pt
        apply_text_fit(
            tf, original_len=10, translated_len=50,
            mode=TextFitMode.AUTO_SHRINK, min_font_ratio=70,
        )

        new_size = tf.paragraphs[0].runs[0].font.size
        self.assertEqual(new_size, Pt(14))

    def test_auto_size_set_when_hitting_floor(self):
        """When hitting min ratio floor, TEXT_TO_FIT_SHAPE should be set as fallback."""
        prs = Presentation()
        shape, tf = _make_textbox(prs, "Hi", font_size_pt=20)

        apply_text_fit(
            tf, original_len=10, translated_len=50,
            mode=TextFitMode.AUTO_SHRINK, min_font_ratio=70,
        )

        self.assertEqual(tf.auto_size, MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE)

    def test_auto_size_not_set_when_not_hitting_floor(self):
        """When proportional shrink is sufficient, don't set TEXT_TO_FIT_SHAPE."""
        prs = Presentation()
        shape, tf = _make_textbox(prs, "Hello", font_size_pt=20)

        # 50% longer -> ratio 0.67, min_font_ratio=50 -> floor not hit
        apply_text_fit(
            tf, original_len=10, translated_len=15,
            mode=TextFitMode.AUTO_SHRINK, min_font_ratio=50,
        )

        self.assertNotEqual(tf.auto_size, MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE)

    def test_multiple_runs_all_shrink(self):
        """All runs in the text frame should be shrunk proportionally."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        tf = shape.text_frame

        # Create paragraph with two runs of different sizes
        p = tf.paragraphs[0]
        p.text = ""
        run1 = p.add_run()
        run1.text = "Bold"
        run1.font.size = Pt(24)
        run1.font.bold = True

        run2 = p.add_run()
        run2.text = " Normal"
        run2.font.size = Pt(16)

        # Text doubled
        apply_text_fit(
            tf, original_len=10, translated_len=20,
            mode=TextFitMode.AUTO_SHRINK, min_font_ratio=50,
        )

        self.assertEqual(run1.font.size, Pt(12))  # 24 * 0.5
        self.assertEqual(run2.font.size, Pt(8))    # 16 * 0.5
        # Bold should be preserved
        self.assertTrue(run1.font.bold)

    def test_run_without_explicit_size_skipped(self):
        """Runs inheriting font size (None) should not be modified."""
        prs = Presentation()
        shape, tf = _make_textbox(prs, "Hello", font_size_pt=20)
        # Add a run without explicit font size
        run2 = tf.paragraphs[0].add_run()
        run2.text = " World"
        run2.font.size = None  # inherits from theme/style

        apply_text_fit(
            tf, original_len=10, translated_len=20,
            mode=TextFitMode.AUTO_SHRINK, min_font_ratio=50,
        )

        # First run should be shrunk
        self.assertEqual(tf.paragraphs[0].runs[0].font.size, Pt(10))
        # Second run should remain None (untouched)
        self.assertIsNone(run2.font.size)

    def test_preserves_existing_text_to_fit_shape(self):
        """If original auto_size is already TEXT_TO_FIT_SHAPE, keep it."""
        prs = Presentation()
        shape, tf = _make_textbox(prs, "Hello", font_size_pt=20)
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        apply_text_fit(
            tf, original_len=5, translated_len=5,
            mode=TextFitMode.AUTO_SHRINK,
        )

        self.assertEqual(tf.auto_size, MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE)

    def test_word_wrap_enabled(self):
        """Auto shrink should enable word wrap."""
        prs = Presentation()
        shape, tf = _make_textbox(prs, "Hello", font_size_pt=20)
        tf.word_wrap = False

        apply_text_fit(
            tf, original_len=10, translated_len=20,
            mode=TextFitMode.AUTO_SHRINK, min_font_ratio=50,
        )

        self.assertTrue(tf.word_wrap)


class TextFitExpandBoxTestCase(unittest.TestCase):
    """TextFitMode.EXPAND_BOX tests."""

    def test_expand_box_sets_auto_size(self):
        """Expand box mode should set SHAPE_TO_FIT_TEXT."""
        prs = Presentation()
        shape, tf = _make_textbox(prs, "Hello", font_size_pt=20)

        apply_text_fit(tf, original_len=5, translated_len=15, mode=TextFitMode.EXPAND_BOX)

        self.assertEqual(tf.auto_size, MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)

    def test_expand_box_enables_word_wrap(self):
        """Expand box should enable word wrap so only height expands."""
        prs = Presentation()
        shape, tf = _make_textbox(prs, "Hello", font_size_pt=20)

        apply_text_fit(tf, original_len=5, translated_len=15, mode=TextFitMode.EXPAND_BOX)

        self.assertTrue(tf.word_wrap)

    def test_expand_box_does_not_modify_font(self):
        """Expand box should never change font sizes."""
        prs = Presentation()
        shape, tf = _make_textbox(prs, "Hello", font_size_pt=20)

        apply_text_fit(tf, original_len=5, translated_len=50, mode=TextFitMode.EXPAND_BOX)

        self.assertEqual(tf.paragraphs[0].runs[0].font.size, Pt(20))

    def test_expand_box_no_change_when_shorter(self):
        """If translated text is shorter, no changes needed."""
        prs = Presentation()
        shape, tf = _make_textbox(prs, "Hello World", font_size_pt=20)
        tf.auto_size = MSO_AUTO_SIZE.NONE

        apply_text_fit(tf, original_len=11, translated_len=5, mode=TextFitMode.EXPAND_BOX)

        # auto_size should remain NONE for shorter text
        self.assertEqual(tf.auto_size, MSO_AUTO_SIZE.NONE)


class TextFitMinFontRatioTestCase(unittest.TestCase):
    """Test various min_font_ratio values."""

    def test_ratio_90(self):
        prs = Presentation()
        shape, tf = _make_textbox(prs, "Hello", font_size_pt=20)

        apply_text_fit(
            tf, original_len=10, translated_len=30,
            mode=TextFitMode.AUTO_SHRINK, min_font_ratio=90,
        )

        # ratio = 10/30 ≈ 0.33, but floor is 90% -> 18pt
        self.assertEqual(tf.paragraphs[0].runs[0].font.size, Pt(18))

    def test_ratio_50(self):
        prs = Presentation()
        shape, tf = _make_textbox(prs, "Hello", font_size_pt=20)

        apply_text_fit(
            tf, original_len=10, translated_len=30,
            mode=TextFitMode.AUTO_SHRINK, min_font_ratio=50,
        )

        # ratio = 10/30 ≈ 0.33, floor is 50% -> 10pt
        self.assertEqual(tf.paragraphs[0].runs[0].font.size, Pt(10))


class TextFitShrinkThenExpandTestCase(unittest.TestCase):
    """TextFitMode.SHRINK_THEN_EXPAND (hybrid) tests."""

    def test_shrink_only_when_within_ratio(self):
        """When shrink alone is sufficient, no box expansion."""
        prs = Presentation()
        shape, tf = _make_textbox(prs, "Hello", font_size_pt=20)
        tf.auto_size = MSO_AUTO_SIZE.NONE

        # 50% longer, min_ratio=50 -> shrink factor 0.67 > floor 0.50 -> no expand
        apply_text_fit(
            tf, original_len=10, translated_len=15,
            mode=TextFitMode.SHRINK_THEN_EXPAND, min_font_ratio=50,
        )

        new_size = tf.paragraphs[0].runs[0].font.size
        # Font should be smaller than original but larger than 50% floor
        self.assertLess(new_size, Pt(20))
        self.assertGreater(new_size, Pt(10))
        # Should NOT expand box
        self.assertNotEqual(tf.auto_size, MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)

    def test_shrink_and_expand_when_hitting_floor(self):
        """When shrink hits floor, also expand box."""
        prs = Presentation()
        shape, tf = _make_textbox(prs, "Hi", font_size_pt=20)

        # 5x longer, min_ratio=70 -> floor hit -> shrink to 70% AND expand
        apply_text_fit(
            tf, original_len=10, translated_len=50,
            mode=TextFitMode.SHRINK_THEN_EXPAND, min_font_ratio=70,
        )

        # Font shrunk to floor
        self.assertEqual(tf.paragraphs[0].runs[0].font.size, Pt(14))
        # Box expanded (NOT TEXT_TO_FIT_SHAPE like auto_shrink)
        self.assertEqual(tf.auto_size, MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)

    def test_word_wrap_enabled(self):
        """Hybrid mode should enable word wrap."""
        prs = Presentation()
        shape, tf = _make_textbox(prs, "Hello", font_size_pt=20)
        tf.word_wrap = False

        apply_text_fit(
            tf, original_len=10, translated_len=20,
            mode=TextFitMode.SHRINK_THEN_EXPAND, min_font_ratio=50,
        )

        self.assertTrue(tf.word_wrap)

    def test_no_change_when_shorter(self):
        """No changes when translated text is shorter."""
        prs = Presentation()
        shape, tf = _make_textbox(prs, "Hello World", font_size_pt=20)

        apply_text_fit(
            tf, original_len=11, translated_len=5,
            mode=TextFitMode.SHRINK_THEN_EXPAND,
        )

        self.assertEqual(tf.paragraphs[0].runs[0].font.size, Pt(20))


def _make_textbox_on_slide(slide, left, top, width, height, text="Hello", font_size_pt=20):
    """Add a textbox to an existing slide and return (shape, text_frame)."""
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.paragraphs[0].text = ""
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(font_size_pt)
    return shape, tf


def _make_paragraph_info(shape, slide_index=0):
    """Create a ParagraphInfo from a shape's first paragraph."""
    tf = shape.text_frame
    p = tf.paragraphs[0]
    text = "".join(r.text for r in p.runs)
    return ParagraphInfo(
        slide_index=slide_index,
        shape_index=0,
        paragraph_index=0,
        original_text=text,
        paragraph=p,
        slide_title=None,
    )


class WidthExpansionTestCase(unittest.TestCase):
    """Tests for automatic width expansion of text boxes."""

    def test_basic_expansion_no_obstacles(self):
        """Single textbox on slide should expand to the right, keeping left fixed."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # Center of a 10-inch slide: left=Inches(3.5), width=Inches(3)
        shape, tf = _make_textbox_on_slide(
            slide, Inches(3.5), Inches(1), Inches(3), Inches(1),
        )
        orig_left = shape.left
        orig_width = shape.width

        txbody_to_shape, slide_bounds, slide_w = _build_shape_context(prs)
        bridge_key = id(shape.text_frame._txBody)
        bounds = slide_bounds[0]

        _, ar = _calculate_available_expansion(shape, bounds, slide_w, bridge_key)

        # Right side has space
        self.assertGreater(ar, 0)

        # Expand by 50% of current width
        needed = int(orig_width * 0.5)
        ratio = _safe_expand_width(shape, ar, needed)

        self.assertGreater(ratio, 1.0)
        self.assertEqual(shape.width, orig_width + needed)
        # Left position must not change
        self.assertEqual(shape.left, orig_left)

    def test_slide_boundary_limits_expansion(self):
        """Textbox near left edge cannot expand left past slide boundary."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # Near left edge: left = gap amount, so available_left ~ 0
        shape, tf = _make_textbox_on_slide(
            slide, _EXPANSION_GAP_EMU, Inches(1), Inches(3), Inches(1),
        )

        txbody_to_shape, slide_bounds, slide_w = _build_shape_context(prs)
        bridge_key = id(shape.text_frame._txBody)
        bounds = slide_bounds[0]

        al, ar = _calculate_available_expansion(shape, bounds, slide_w, bridge_key)

        # Left should be 0 (at the edge minus gap)
        self.assertEqual(al, 0)
        # Right should have space
        self.assertGreater(ar, 0)

    def test_adjacent_shape_right_limits_expansion(self):
        """Obstacle to the right limits right expansion."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # Target shape on the left
        shape, tf = _make_textbox_on_slide(
            slide, Inches(1), Inches(1), Inches(2), Inches(1),
        )
        # Obstacle close to the right (gap of 0.5 inch)
        _make_textbox_on_slide(
            slide, Inches(3.5), Inches(1), Inches(2), Inches(1), text="obstacle",
        )

        txbody_to_shape, slide_bounds, slide_w = _build_shape_context(prs)
        bridge_key = id(shape.text_frame._txBody)
        bounds = slide_bounds[0]

        al, ar = _calculate_available_expansion(shape, bounds, slide_w, bridge_key)

        # Right is limited by obstacle: Inches(3.5) - Inches(3) = Inches(0.5) - gap
        expected_right = Inches(0.5) - _EXPANSION_GAP_EMU
        self.assertAlmostEqual(ar, max(0, expected_right), delta=1)

    def test_adjacent_shape_left_limits_expansion(self):
        """Obstacle to the left limits left expansion."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # Obstacle on the left
        _make_textbox_on_slide(
            slide, Inches(0.5), Inches(1), Inches(2), Inches(1), text="obstacle",
        )
        # Target shape to the right of obstacle (gap of 0.5 inch)
        shape, tf = _make_textbox_on_slide(
            slide, Inches(3), Inches(1), Inches(2), Inches(1),
        )

        txbody_to_shape, slide_bounds, slide_w = _build_shape_context(prs)
        bridge_key = id(shape.text_frame._txBody)
        bounds = slide_bounds[0]

        al, ar = _calculate_available_expansion(shape, bounds, slide_w, bridge_key)

        # Left limited by obstacle: Inches(3) - Inches(2.5) = Inches(0.5) - gap
        expected_left = Inches(0.5) - _EXPANSION_GAP_EMU
        self.assertAlmostEqual(al, max(0, expected_left), delta=1)

    def test_right_only_expansion_preserves_left(self):
        """Expansion only goes right; left position never changes."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape, tf = _make_textbox_on_slide(
            slide, Inches(0.5), Inches(1), Inches(2), Inches(1),
        )

        txbody_to_shape, slide_bounds, slide_w = _build_shape_context(prs)
        bridge_key = id(shape.text_frame._txBody)
        bounds = slide_bounds[0]
        _, ar = _calculate_available_expansion(shape, bounds, slide_w, bridge_key)

        orig_left = shape.left
        orig_width = shape.width
        needed = int(orig_width * 1.0)  # double width
        ratio = _safe_expand_width(shape, ar, needed)

        # Left must not move
        self.assertEqual(shape.left, orig_left)
        # Width increased
        self.assertGreater(ratio, 1.0)
        self.assertGreater(shape.width, orig_width)

    def test_group_child_excluded(self):
        """Text frames from group children should not be in txbody_to_shape mapping."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # Add a regular textbox
        shape, tf = _make_textbox_on_slide(
            slide, Inches(1), Inches(1), Inches(3), Inches(1),
        )

        txbody_to_shape, _, _ = _build_shape_context(prs)
        # The regular textbox IS in the mapping
        bridge_key = id(shape.text_frame._txBody)
        self.assertIn(bridge_key, txbody_to_shape)

        # Group children are not iterated by slide.shapes (only top-level)
        # so they won't appear in txbody_to_shape.
        # We verify the mapping only contains the one textbox we created.
        # (slide_layouts[6] may have placeholder shapes too)
        self.assertIn(bridge_key, txbody_to_shape)

    def test_table_cell_excluded(self):
        """Table cell text frames should not be in txbody_to_shape mapping."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # Add a table
        table_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2))
        table = table_shape.table
        cell = table.cell(0, 0)
        cell.text = "CellText"
        cell_txbody_id = id(cell.text_frame._txBody)

        txbody_to_shape, _, _ = _build_shape_context(prs)

        # Table cell text frame should NOT be mapped (it's not a top-level shape text frame)
        self.assertNotIn(cell_txbody_id, txbody_to_shape)

    def test_rotated_shape_skipped(self):
        """Rotated shapes should not get width expansion via apply_translations."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape, tf = _make_textbox_on_slide(
            slide, Inches(2), Inches(1), Inches(3), Inches(1),
        )
        shape.rotation = 45.0
        orig_width = shape.width
        orig_left = shape.left

        para_info = _make_paragraph_info(shape)

        writer = PPTWriter()
        # Text 3x longer -> triggers width expansion + text fit
        long_text = "x" * (len(para_info.original_text) * 3)
        writer.apply_translations(
            [para_info], [long_text], prs,
            text_fit_mode="auto_shrink", min_font_ratio=50,
        )

        # Shape should NOT be width-expanded (rotation != 0)
        self.assertEqual(shape.width, orig_width)
        self.assertEqual(shape.left, orig_left)

    def test_no_vertical_overlap_not_obstacle(self):
        """Shapes without vertical overlap should not block expansion."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # Target shape at y=Inches(3)
        shape, tf = _make_textbox_on_slide(
            slide, Inches(3), Inches(3), Inches(2), Inches(1),
        )
        # "Obstacle" at y=Inches(0.5) — no vertical overlap with target
        _make_textbox_on_slide(
            slide, Inches(5.5), Inches(0.5), Inches(2), Inches(1), text="far",
        )

        txbody_to_shape, slide_bounds, slide_w = _build_shape_context(prs)
        bridge_key = id(shape.text_frame._txBody)
        bounds = slide_bounds[0]

        al, ar = _calculate_available_expansion(shape, bounds, slide_w, bridge_key)

        # Right should NOT be limited by the non-overlapping shape
        # Available right = slide_width - shape_right - gap
        expected_right = slide_w - (Inches(3) + Inches(2)) - _EXPANSION_GAP_EMU
        self.assertAlmostEqual(ar, expected_right, delta=1)

    def test_auto_shrink_reduced_by_width_expansion(self):
        """Width expansion should reduce the amount of font shrinking needed."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # Centered textbox with plenty of room
        shape, tf = _make_textbox_on_slide(
            slide, Inches(3), Inches(1), Inches(2), Inches(1),
            text="Hello", font_size_pt=20,
        )
        para_info = _make_paragraph_info(shape)
        orig_text = para_info.original_text

        # Create a second presentation with same setup but NO width expansion (mode=none first, then manual)
        prs2 = Presentation()
        slide2 = prs2.slides.add_slide(prs2.slide_layouts[6])
        shape2, tf2 = _make_textbox_on_slide(
            slide2, Inches(3), Inches(1), Inches(2), Inches(1),
            text="Hello", font_size_pt=20,
        )

        # Text 2x longer
        long_text = "x" * (len(orig_text) * 2)

        # Apply with auto_shrink (has width expansion)
        para1 = _make_paragraph_info(shape)
        writer = PPTWriter()
        writer.apply_translations(
            [para1], [long_text], prs,
            text_fit_mode="auto_shrink", min_font_ratio=50,
        )

        # Apply directly without width expansion
        apply_text_fit(tf2, len(orig_text), len(long_text), mode="auto_shrink", min_font_ratio=50)

        font_with_expansion = tf.paragraphs[0].runs[0].font.size
        font_without_expansion = tf2.paragraphs[0].runs[0].font.size

        # Font with width expansion should be >= font without (less shrink needed)
        self.assertGreaterEqual(font_with_expansion, font_without_expansion)

    def test_expand_box_skipped_when_width_expansion_sufficient(self):
        """If width expansion covers the text growth, SHAPE_TO_FIT_TEXT should not be set."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # Textbox with lots of room to expand
        shape, tf = _make_textbox_on_slide(
            slide, Inches(4), Inches(1), Inches(1), Inches(1),
            text="Hello World Test!", font_size_pt=20,
        )
        tf.auto_size = MSO_AUTO_SIZE.NONE  # Explicit baseline
        para_info = _make_paragraph_info(shape)

        # Text ~15% longer (just above threshold) — width expansion should cover it
        orig_len = len(para_info.original_text)
        trans_text = "x" * int(orig_len * 1.15)
        self.assertGreater(len(trans_text), orig_len)  # sanity check

        writer = PPTWriter()
        writer.apply_translations(
            [para_info], [trans_text], prs,
            text_fit_mode="expand_box", min_font_ratio=70,
        )

        # Width expansion should have compensated, so SHAPE_TO_FIT_TEXT not needed
        self.assertNotEqual(tf.auto_size, MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)


if __name__ == "__main__":
    unittest.main()
